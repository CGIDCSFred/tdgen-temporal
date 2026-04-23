"""
TDGen Temporal — Internal Stakeholder Presentation
CGI Standard Presentation Template (June 2024)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------
CGI_RED = RGBColor(0xE4, 0x00, 0x2B)
CGI_DARK_RED = RGBColor(0x8C, 0x00, 0x15)
CGI_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
CGI_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
CGI_BODY = RGBColor(0x33, 0x33, 0x33)
CGI_GREY = RGBColor(0x6E, 0x6E, 0x6E)
CGI_LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
CGI_MID_GREY = RGBColor(0xCC, 0xCC, 0xCC)
CGI_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CGI_BLUE = RGBColor(0x00, 0x5B, 0x99)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT = "Calibri"
MARGIN_L = Inches(0.55)
CONTENT_W = Inches(12.23)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def _apply_gradient(shape, stops, angle_deg=90.0):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for tag in [
        qn("a:solidFill"),
        qn("a:gradFill"),
        qn("a:noFill"),
        qn("a:blipFill"),
        qn("a:pattFill"),
        qn("a:grpFill"),
    ]:
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    gs_xml = "".join(
        f'<a:gs pos="{int(p * 100000)}"><a:srgbClr val="{c}"/></a:gs>' for p, c in stops
    )
    ang = int(angle_deg * 60000)
    grad_xml = (
        f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f"<a:gsLst>{gs_xml}</a:gsLst>"
        f'<a:lin ang="{ang}" scaled="0"/>'
        f"</a:gradFill>"
    )
    grad_el = parse_xml(grad_xml)
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.addnext(grad_el)
    else:
        spPr.append(grad_el)
    shape.line.fill.background()


def rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def grad_rect(slide, left, top, width, height, stops, angle_deg=90.0):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = CGI_RED
    _apply_gradient(s, stops, angle_deg)
    return s


def _h_seg(slide, x1, x2, y):
    lft = min(x1, x2)
    wid = abs(x2 - x1)
    if wid > 0:
        LINE_W = Inches(0.020)
        rect(slide, lft, y - LINE_W / 2, wid, LINE_W, RGBColor(0x88, 0x88, 0x88))


def _v_seg(slide, y1, y2, x):
    top = min(y1, y2)
    hgt = abs(y2 - y1)
    if hgt > 0:
        LINE_W = Inches(0.020)
        rect(slide, x - LINE_W / 2, top, LINE_W, hgt, RGBColor(0x88, 0x88, 0x88))


def txt(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=16,
    bold=False,
    color=CGI_BODY,
    align=PP_ALIGN.LEFT,
    wrap=True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return box


def multi_para(slide, lines, left, top, width, height, space_pt=4.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        from pptx.util import Pt as _Pt

        p.space_before = _Pt(space_pt if i > 0 else 0)
        run = p.add_run()
        run.text = str(text)
        run.font.size = _Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def cornerstone(slide):
    grad_rect(
        slide,
        Inches(6.80),
        Inches(0),
        Inches(6.53),
        Inches(3.90),
        [(0.0, "FF8C00"), (1.0, "E4002B")],
        135,
    )
    grad_rect(
        slide,
        Inches(10.00),
        Inches(3.90),
        Inches(3.33),
        Inches(3.60),
        [(0.0, "E4002B"), (1.0, "8C0015")],
        90,
    )


def add_footer(slide, page_num):
    y = SLIDE_H - Inches(0.40)
    h = Inches(0.32)
    txt(slide, "(c) 2025 CGI Inc.", MARGIN_L, y, Inches(3.5), h, size=10, color=CGI_GREY)
    txt(
        slide,
        "Internal",
        Inches(5.5),
        y,
        Inches(2.33),
        h,
        size=10,
        color=CGI_GREY,
        align=PP_ALIGN.CENTER,
    )
    txt(
        slide,
        str(page_num),
        SLIDE_W - Inches(1.3),
        y,
        Inches(0.9),
        h,
        size=10,
        color=CGI_GREY,
        align=PP_ALIGN.RIGHT,
    )


def slide_header(slide, title, subtitle=None):
    txt(
        slide,
        title,
        MARGIN_L,
        Inches(0.18),
        CONTENT_W,
        Inches(0.68),
        size=26,
        bold=False,
        color=CGI_BLACK,
    )
    if subtitle:
        rect(slide, MARGIN_L, Inches(0.93), CONTENT_W, Inches(0.022), CGI_MID_GREY)
        txt(
            slide,
            subtitle,
            MARGIN_L,
            Inches(0.97),
            CONTENT_W,
            Inches(0.34),
            size=13,
            color=CGI_GREY,
        )
        rect(slide, MARGIN_L, Inches(1.38), CONTENT_W, Inches(0.022), CGI_MID_GREY)
    else:
        rect(slide, MARGIN_L, Inches(0.93), CONTENT_W, Inches(0.022), CGI_MID_GREY)


# ---------------------------------------------------------------------------
# SLIDE 1  -  Cover
# ---------------------------------------------------------------------------


def slide_01_cover(prs):
    slide = blank_slide(prs)
    cornerstone(slide)

    # Left content panel
    rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, CGI_RED)

    txt(
        slide,
        "TDGen Temporal",
        Inches(0.55),
        Inches(1.40),
        Inches(6.0),
        Inches(0.90),
        size=42,
        bold=True,
        color=CGI_BLACK,
    )
    txt(
        slide,
        "Day-Over-Day Synthetic Test Data Generation",
        Inches(0.55),
        Inches(2.38),
        Inches(6.0),
        Inches(0.55),
        size=22,
        bold=False,
        color=CGI_BODY,
    )

    rect(slide, Inches(0.55), Inches(3.04), Inches(5.80), Inches(0.04), CGI_RED)

    multi_para(
        slide,
        [
            ("TSYS TS2  |  Credit Card Fraud Lifecycle", 15, False, CGI_BODY),
            ("TD Bank  |  CGI Consulting Engagement", 15, False, CGI_BODY),
            ("April 2026", 14, False, CGI_GREY),
        ],
        Inches(0.55),
        Inches(3.18),
        Inches(6.0),
        Inches(1.10),
        space_pt=5,
    )

    add_footer(slide, 1)


# ---------------------------------------------------------------------------
# SLIDE 2  -  Introduction / Overview
# ---------------------------------------------------------------------------


def slide_02_what(prs):
    slide = blank_slide(prs)

    # Title + rule
    txt(
        slide,
        "Temporal Test Data Generator (TDGen Temporal)",
        MARGIN_L,
        Inches(0.18),
        CONTENT_W,
        Inches(0.68),
        size=26,
        bold=False,
        color=CGI_BLACK,
    )
    rect(slide, MARGIN_L, Inches(0.93), CONTENT_W, Inches(0.022), CGI_MID_GREY)

    LEFT_W = Inches(7.90)
    RIGHT_X = MARGIN_L + LEFT_W + Inches(0.22)
    RIGHT_W = SLIDE_W - RIGHT_X - Inches(0.18)

    # ------------------------------------------------------------------ #
    # PURPOSE
    # ------------------------------------------------------------------ #
    Y = Inches(1.04)
    txt(slide, "PURPOSE", MARGIN_L, Y, LEFT_W, Inches(0.24), size=10, bold=True, color=CGI_RED)
    Y += Inches(0.27)

    multi_para(
        slide,
        [
            (
                "\u2022  Enable CGI QE teams at TD Bank to generate high-fidelity synthetic test data on demand "
                "\u2014 directly from current schemas \u2014 without extracting data from source platforms or "
                "coordinating with source system teams.",
                12,
                False,
                CGI_BODY,
            ),
            (
                "\u2022  TDGen Temporal is a standalone Python CLI with persistent SQLite state. "
                "No dependency on Claude Code or any AI service at runtime. "
                "The MCP server is an optional interaction layer.",
                12,
                False,
                CGI_BODY,
            ),
        ],
        MARGIN_L,
        Y,
        LEFT_W,
        Inches(1.12),
        space_pt=5,
    )
    Y += Inches(1.18)

    # ------------------------------------------------------------------ #
    # APPROACH  (5-step pipeline)
    # ------------------------------------------------------------------ #
    txt(slide, "APPROACH", MARGIN_L, Y, LEFT_W, Inches(0.24), size=10, bold=True, color=CGI_RED)
    Y += Inches(0.27)

    steps = [
        (RGBColor(0x1A, 0x3C, 0x5E), "DB", "Day 0 seed", "500 accounts\n+ cards"),
        (RGBColor(0x1B, 0x5E, 0x3B), "\u2699", "Daily advance", "State machines\n+ transactions"),
        (RGBColor(0x7A, 0x5A, 0x00), "\u25d4", "Backfill", "Historical\ndate ranges"),
        (RGBColor(0x5E, 0x1A, 0x6E), "\u25a6", "Delta output", "INSERT + UPDATE\nper day"),
        (RGBColor(0x1A, 0x3C, 0x6E), "\u29bf", "MCP server", "Optional\nlayer"),
    ]

    STEP_W = Inches(1.48)
    STEP_H = Inches(1.76)
    CHEV_W = Inches(0.10)

    sx = MARGIN_L
    for idx, (color, icon, title, sub) in enumerate(steps):
        rect(slide, sx, Y, STEP_W, STEP_H, color)
        ICON_S = Inches(0.54)
        icon_x = sx + (STEP_W - ICON_S) / 2
        rect(slide, int(icon_x), Y + Inches(0.12), ICON_S, ICON_S, RGBColor(0xFF, 0xFF, 0xFF))
        txt(
            slide,
            icon,
            int(icon_x),
            Y + Inches(0.12),
            ICON_S,
            ICON_S,
            size=20,
            bold=True,
            color=color,
            align=PP_ALIGN.CENTER,
        )
        txt(
            slide,
            title,
            sx + Inches(0.08),
            Y + Inches(0.82),
            STEP_W - Inches(0.16),
            Inches(0.36),
            size=12,
            bold=True,
            color=CGI_WHITE,
        )
        txt(
            slide,
            sub,
            sx + Inches(0.08),
            Y + Inches(1.20),
            STEP_W - Inches(0.16),
            Inches(0.50),
            size=10,
            color=RGBColor(0xBB, 0xBB, 0xBB),
        )
        sx += STEP_W
        if idx < 4:
            txt(
                slide,
                "\u203a",
                sx + Inches(0.01),
                Y + STEP_H / 2 - Inches(0.18),
                CHEV_W + Inches(0.02),
                Inches(0.36),
                size=18,
                bold=True,
                color=CGI_GREY,
                align=PP_ALIGN.CENTER,
            )
            sx += CHEV_W

    Y += STEP_H + Inches(0.22)

    # ------------------------------------------------------------------ #
    # KEY OUTCOMES
    # ------------------------------------------------------------------ #
    txt(slide, "KEY OUTCOMES", MARGIN_L, Y, LEFT_W, Inches(0.24), size=10, bold=True, color=CGI_RED)
    Y += Inches(0.27)

    outcomes = [
        ("Days \u2192 seconds", "full history backfill"),
        ("Zero source data", "no cross-team coordination"),
        ("Standalone pipeline", "no runtime LLM"),
        ("Compliance safe", "100% synthetic, no PDIT"),
    ]

    OUT_W = (LEFT_W - 3 * Inches(0.09)) // 4
    OUT_H = Inches(0.84)
    ox = MARGIN_L
    for main, sub in outcomes:
        rect(slide, ox, Y, OUT_W, OUT_H, CGI_BLACK)
        txt(
            slide,
            main,
            ox + Inches(0.10),
            Y + Inches(0.07),
            OUT_W - Inches(0.16),
            Inches(0.38),
            size=12,
            bold=True,
            color=CGI_WHITE,
        )
        txt(
            slide,
            sub,
            ox + Inches(0.10),
            Y + Inches(0.48),
            OUT_W - Inches(0.16),
            Inches(0.30),
            size=10,
            color=CGI_GREY,
        )
        ox += OUT_W + Inches(0.09)

    # ------------------------------------------------------------------ #
    # TOOL STACK CARD  (right panel)
    # ------------------------------------------------------------------ #
    CARD_Y = Inches(1.00)
    CARD_H = SLIDE_H - CARD_Y - Inches(0.48)
    rect(slide, RIGHT_X, CARD_Y, RIGHT_W, CARD_H, CGI_LIGHT_GREY)
    rect(slide, RIGHT_X, CARD_Y, Inches(0.05), CARD_H, CGI_RED)

    txt(
        slide,
        "Tool Stack:",
        RIGHT_X + Inches(0.14),
        CARD_Y + Inches(0.14),
        RIGHT_W - Inches(0.22),
        Inches(0.38),
        size=14,
        bold=True,
        color=CGI_BLACK,
    )

    tools = [
        "Agentic Dev \u2013 Claude Code CLI w/Sonnet 4.6",
        "IDE \u2013 Visual Studio Code w/Claude Extension",
        "Core language \u2013 Python 3.11+",
        "MCP server tool -- FastMCP",
        "Schema ingestion -- openpyxl",
        "Synthetic data \u2013 Faker",
        "Data Validation \u2013 Pandera",
        "DBMS \u2013 SQLite",
        "ETL \u2013 pandas",
        "Config \u2013 PyYAML",
        "CI/CD \u2013 GitHub Actions",
    ]

    multi_para(
        slide,
        [("\u2022  " + t, 11, False, CGI_BODY) for t in tools],
        RIGHT_X + Inches(0.14),
        CARD_Y + Inches(0.58),
        RIGHT_W - Inches(0.22),
        CARD_H - Inches(0.70),
        space_pt=3.5,
    )

    add_footer(slide, 2)


# ---------------------------------------------------------------------------
# SLIDE 3  -  POC Scope
# ---------------------------------------------------------------------------


def slide_03_scope(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "POC Scope  \u2014  Data Model & Simulation Engine",
        "TSYS TS2 approximation derived from publicly available sources  |  Claude-built model",
    )

    PANEL_TOP = Inches(1.56)
    PANEL_H = Inches(3.30)
    LEFT_W = Inches(5.80)
    GAP_MID = Inches(0.30)
    RIGHT_X = MARGIN_L + LEFT_W + GAP_MID
    RIGHT_W = CONTENT_W - LEFT_W - GAP_MID

    # ── Left panel: Data Model ──────────────────────────────────────────────
    rect(slide, MARGIN_L, PANEL_TOP, LEFT_W, PANEL_H, CGI_LIGHT_GREY)
    rect(slide, MARGIN_L, PANEL_TOP, Inches(0.06), PANEL_H, CGI_RED)

    txt(
        slide,
        "DATA MODEL",
        MARGIN_L + Inches(0.14),
        PANEL_TOP + Inches(0.12),
        LEFT_W - Inches(0.20),
        Inches(0.22),
        size=10,
        bold=True,
        color=CGI_RED,
    )
    txt(
        slide,
        "TSYS TS2  \u2014  publicly available approximation, built by Claude",
        MARGIN_L + Inches(0.14),
        PANEL_TOP + Inches(0.38),
        LEFT_W - Inches(0.20),
        Inches(0.28),
        size=12,
        bold=False,
        color=CGI_BLACK,
    )

    # Five stat tiles
    stats = [
        ("31", "Tables"),
        ("15", "Entities"),
        ("16", "Lookups"),
        ("285", "Columns"),
        ("47", "Relationships"),
    ]
    INNER_W = LEFT_W - Inches(0.28)
    TILE_GAP = Inches(0.06)
    TILE_W = (INNER_W - 4 * TILE_GAP) // 5
    TILE_H = Inches(1.72)
    tx = MARGIN_L + Inches(0.14)
    ty = PANEL_TOP + Inches(0.78)

    for num, label in stats:
        rect(slide, tx, ty, TILE_W, TILE_H, CGI_BLACK)
        txt(
            slide,
            num,
            tx,
            ty + Inches(0.18),
            TILE_W,
            Inches(0.90),
            size=32,
            bold=True,
            color=CGI_WHITE,
            align=PP_ALIGN.CENTER,
        )
        txt(
            slide,
            label,
            tx,
            ty + Inches(1.14),
            TILE_W,
            Inches(0.44),
            size=9,
            color=RGBColor(0xBB, 0xBB, 0xBB),
            align=PP_ALIGN.CENTER,
        )
        tx += TILE_W + TILE_GAP

    # ── Right panel: Entity State Machines ──────────────────────────────────
    rect(slide, RIGHT_X, PANEL_TOP, RIGHT_W, PANEL_H, CGI_LIGHT_GREY)
    rect(slide, RIGHT_X, PANEL_TOP, Inches(0.06), PANEL_H, CGI_RED)

    txt(
        slide,
        "ENTITY STATE MACHINES",
        RIGHT_X + Inches(0.14),
        PANEL_TOP + Inches(0.12),
        RIGHT_W - Inches(0.20),
        Inches(0.22),
        size=10,
        bold=True,
        color=CGI_RED,
    )
    txt(
        slide,
        "6 entities advance through defined lifecycle states each simulation day",
        RIGHT_X + Inches(0.14),
        PANEL_TOP + Inches(0.38),
        RIGHT_W - Inches(0.20),
        Inches(0.28),
        size=12,
        bold=False,
        color=CGI_BLACK,
    )

    machines = [
        (
            "ACCOUNT",
            "ACTIVE \u2192 DELINQUENT \u2192 CHARGEOFF \u2192 CLOSED",
            RGBColor(0x00, 0x5B, 0x99),
        ),
        ("CARD", "ACTIVE \u2192 BLOCKED \u2192 EXPIRED \u2192 CLOSED", RGBColor(0x00, 0x6E, 0x51)),
        (
            "DISPUTE",
            "OPEN \u2192 INVESTIGATING \u2192 RESOLVED \u2192 CLOSED",
            RGBColor(0x1A, 0x7A, 0x3C),
        ),
        ("FRAUD ALERT", "OPEN \u2192 UNDER_REVIEW \u2192 CONFIRMED | FALSE_POS", CGI_ORANGE),
        (
            "CHARGEBACK",
            "FIRST_CB \u2192 REPRESENTMENT \u2192 WON | LOST",
            RGBColor(0x5E, 0x1A, 0x6E),
        ),
        (
            "COLLECTION",
            "B1 \u2192 B2 \u2192 B3 \u2192 B4 \u2192 CHARGEOFF",
            RGBColor(0x1A, 0x3C, 0x5E),
        ),
    ]

    SM_H = Inches(0.38)
    SM_GAP = Inches(0.04)
    sy = PANEL_TOP + Inches(0.78)
    NAME_W = Inches(1.30)

    for name, states, color in machines:
        rect(slide, RIGHT_X + Inches(0.14), sy, RIGHT_W - Inches(0.28), SM_H, CGI_WHITE)
        rect(slide, RIGHT_X + Inches(0.14), sy, Inches(0.04), SM_H, color)
        txt(
            slide,
            name,
            RIGHT_X + Inches(0.24),
            sy + Inches(0.08),
            NAME_W,
            Inches(0.24),
            size=11,
            bold=True,
            color=color,
        )
        txt(
            slide,
            states,
            RIGHT_X + Inches(0.24) + NAME_W,
            sy + Inches(0.10),
            RIGHT_W - Inches(0.28) - NAME_W - Inches(0.10),
            Inches(0.22),
            size=10,
            color=CGI_BODY,
        )
        sy += SM_H + SM_GAP

    # ── Bottom: Daily Simulation Loop ───────────────────────────────────────
    BOTTOM_Y = PANEL_TOP + PANEL_H + Inches(0.18)
    BOTTOM_H = SLIDE_H - BOTTOM_Y - Inches(0.48)

    rect(slide, MARGIN_L, BOTTOM_Y, CONTENT_W, BOTTOM_H, CGI_LIGHT_GREY)
    rect(slide, MARGIN_L, BOTTOM_Y, Inches(0.06), BOTTOM_H, CGI_RED)

    txt(
        slide,
        "DAILY SIMULATION LOOP",
        MARGIN_L + Inches(0.14),
        BOTTOM_Y + Inches(0.12),
        CONTENT_W - Inches(0.22),
        Inches(0.22),
        size=10,
        bold=True,
        color=CGI_RED,
    )

    loop_text = (
        "Each day, the simulator loads all active accounts and advances every entity through its lifecycle "
        "\u2014 updating account and card states, generating transactions, statements, and scores, then "
        "progressing any open disputes, fraud alerts, chargebacks, and collection cases "
        "\u2014 before persisting all changes to the database and exporting the day\u2019s delta files."
    )
    txt(
        slide,
        loop_text,
        MARGIN_L + Inches(0.14),
        BOTTOM_Y + Inches(0.40),
        CONTENT_W - Inches(0.22),
        BOTTOM_H - Inches(0.50),
        size=12,
        color=CGI_BODY,
    )

    add_footer(slide, 3)


# ---------------------------------------------------------------------------
# SLIDE 4  -  How It Works  (pipeline flow)
# ---------------------------------------------------------------------------


def slide_03_how(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "How It Works",
        "13-step daily execution sequence — all entities advance in dependency order",
    )

    steps = [
        ("1", "Load accounts", "Active accounts + temporal state from SQLite"),
        (
            "2",
            "Account state machine",
            "Balance updates, delinquency aging, ACTIVE->DELINQUENT->CHARGEOFF",
        ),
        ("3", "Card state machine", "Expiry checks, BLOCKED cards, replacement issuance"),
        ("4", "Score refresh", "Monthly risk score drift based on account health"),
        ("5", "Generate transactions", "~1.8 txns/account/day  (purchases, payments, fees)"),
        ("6", "Generate authorizations", "One auth per transaction + declined-only auths"),
        ("7", "Generate statements", "Triggered when run_date.day == account.cycle_day"),
        ("8", "Open disputes", "Fraction of purchases become disputes (FRAUD, DUPLICATE, etc.)"),
        ("9", "Open fraud alerts", "Flagged transactions trigger FRAUD_ALERT records"),
        ("10", "Advance disputes", "OPEN -> INVESTIGATING -> RESOLVED -> CLOSED"),
        ("11", "Advance fraud alerts", "OPEN -> UNDER_REVIEW -> CONFIRMED | FALSE_POSITIVE"),
        ("12", "Advance chargebacks", "FIRST_CHARGEBACK -> REPRESENTMENT -> WON | LOST"),
        ("13", "Advance collection cases", "Bucket progression B1->B2->B3->B4->CHARGEOFF"),
    ]

    col_w = Inches(5.90)
    row_h = Inches(0.375)
    gap = Inches(0.06)
    top = Inches(1.56)
    left1 = MARGIN_L
    left2 = MARGIN_L + col_w + Inches(0.43)

    for i, (num, heading, detail) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = left1 if col == 0 else left2
        y = top + row * (row_h + gap)
        bg = CGI_LIGHT_GREY if row % 2 == 0 else CGI_WHITE

        rect(slide, x, y, col_w, row_h, bg)
        rect(slide, x, y, Inches(0.04), row_h, CGI_RED)

        txt(
            slide,
            num,
            x + Inches(0.10),
            y + Inches(0.05),
            Inches(0.30),
            Inches(0.28),
            size=13,
            bold=True,
            color=CGI_RED,
        )
        txt(
            slide,
            heading,
            x + Inches(0.42),
            y + Inches(0.05),
            Inches(1.60),
            Inches(0.28),
            size=12,
            bold=True,
            color=CGI_BLACK,
        )
        txt(
            slide,
            detail,
            x + Inches(2.08),
            y + Inches(0.05),
            col_w - Inches(2.16),
            Inches(0.28),
            size=12,
            color=CGI_BODY,
        )

    add_footer(slide, 4)


# ---------------------------------------------------------------------------
# SLIDE 5  -  State Machines
# ---------------------------------------------------------------------------


def slide_04_state_machines(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Entity State Machines",
        "Six entities advance through defined lifecycle states each day",
    )

    machines = [
        (
            "ACCOUNT",
            ["ACTIVE", "DELINQUENT", "CHARGEOFF", "CLOSED"],
            [
                ("ACTIVE", "DELINQUENT", "Missed payment"),
                ("DELINQUENT", "ACTIVE", "Payment received"),
                ("DELINQUENT", "CHARGEOFF", "180+ days past due"),
                ("CHARGEOFF", "CLOSED", "Recovery complete"),
            ],
        ),
        (
            "DISPUTE",
            ["OPEN", "INVESTIGATING", "RESOLVED", "CLOSED"],
            [
                ("OPEN", "INVESTIGATING", "Auto next-day pickup"),
                ("OPEN", "WITHDRAWN", "Customer cancels (<3 days)"),
                ("INVESTIGATING", "RESOLVED", "7+ days, 70% probability"),
                ("RESOLVED", "CLOSED", "5-day settlement window"),
            ],
        ),
        (
            "FRAUD_ALERT",
            ["OPEN", "UNDER_REVIEW", "CONFIRMED", "FALSE_POSITIVE"],
            [
                ("OPEN", "UNDER_REVIEW", "95% next-day pickup"),
                ("UNDER_REVIEW", "CONFIRMED", "2+ days, 30% rate -> block card"),
                ("UNDER_REVIEW", "FALSE_POSITIVE", "70% rate -> restore card"),
                ("CONFIRMED", "CLOSED", "30-day closure sweep"),
            ],
        ),
    ]

    BOX_W = Inches(3.80)
    BOX_H = Inches(5.10)
    GAP = Inches(0.215)
    TOP = Inches(1.56)
    X = MARGIN_L

    STATE_H = Inches(0.36)
    TRANS_H = Inches(0.32)

    for name, states, transitions in machines:
        rect(slide, X, TOP, BOX_W, BOX_H, CGI_LIGHT_GREY)
        rect(slide, X, TOP, Inches(0.06), BOX_H, CGI_RED)

        # Entity name header
        rect(slide, X, TOP, BOX_W, Inches(0.42), CGI_RED)
        txt(
            slide,
            name,
            X + Inches(0.12),
            TOP + Inches(0.04),
            BOX_W - Inches(0.16),
            Inches(0.36),
            size=15,
            bold=True,
            color=CGI_WHITE,
        )

        # States
        sy = TOP + Inches(0.52)
        txt(
            slide,
            "States",
            X + Inches(0.12),
            sy,
            BOX_W - Inches(0.16),
            Inches(0.24),
            size=11,
            bold=True,
            color=CGI_GREY,
        )
        sy += Inches(0.26)
        for state in states:
            rect(
                slide,
                X + Inches(0.12),
                sy,
                BOX_W - Inches(0.24),
                STATE_H,
                RGBColor(0xFF, 0xFF, 0xFF),
            )
            rect(slide, X + Inches(0.12), sy, Inches(0.04), STATE_H, CGI_BLUE)
            txt(
                slide,
                state,
                X + Inches(0.22),
                sy + Inches(0.04),
                BOX_W - Inches(0.36),
                Inches(0.28),
                size=12,
                bold=False,
                color=CGI_BLACK,
            )
            sy += STATE_H + Inches(0.04)

        # Transitions
        sy += Inches(0.06)
        txt(
            slide,
            "Transitions",
            X + Inches(0.12),
            sy,
            BOX_W - Inches(0.16),
            Inches(0.24),
            size=11,
            bold=True,
            color=CGI_GREY,
        )
        sy += Inches(0.26)
        for fr, to, cond in transitions:
            txt(
                slide,
                f"{fr} -> {to}",
                X + Inches(0.12),
                sy,
                BOX_W - Inches(0.16),
                Inches(0.20),
                size=11,
                bold=True,
                color=CGI_RED,
            )
            sy += Inches(0.22)
            txt(
                slide,
                cond,
                X + Inches(0.18),
                sy,
                BOX_W - Inches(0.22),
                Inches(0.20),
                size=11,
                color=CGI_BODY,
            )
            sy += Inches(0.24)

        X += BOX_W + GAP

    add_footer(slide, 5)


# ---------------------------------------------------------------------------
# SLIDE 5  -  Tech Stack
# ---------------------------------------------------------------------------


def slide_05_stack(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Technology Stack",
        "Standalone Python — no dependency on the original TDGen batch project",
    )

    stack = [
        (
            "Agentic Coding",
            "Claude Code CLI",
            "Claude Sonnet 4.6",
            "Built and iterated entirely via Claude Code in a single session",
        ),
        (
            "IDE",
            "Visual Studio Code",
            "with Claude Code extension",
            "Full codebase visible; hooks into terminal for live testing",
        ),
        (
            "Language",
            "Python 3.11+",
            "",
            "Type hints throughout; dataclasses for all structured data models",
        ),
        (
            "Schema Ingestion",
            "openpyxl",
            "",
            "Reads TSYS TS2 workbook for REF table seed data and column definitions",
        ),
        (
            "Data Generation",
            "Faker (en_CA)",
            "",
            "Canadian locale — SINs, postal codes, phone numbers, names",
        ),
        (
            "State Persistence",
            "SQLite",
            "stdlib sqlite3 + WAL mode",
            "Single state.db file; all entities + temporal state tables",
        ),
        ("ETL / Output", "pandas", "", "DataFrame-to-CSV for delta and snapshot output files"),
        (
            "Data Validation",
            "pandera",
            "replaced Great Expectations",
            "DataFrame schema validation; FK integrity checks per run",
        ),
        (
            "Config",
            "PyYAML",
            "scenario.yaml",
            "All rates, thresholds, and population sizes configurable without code changes",
        ),
        ("CI/CD", "GitHub Actions", "", "Scheduled daily advance + manual dispatch workflow"),
    ]

    row_h = Inches(0.470)
    gap = Inches(0.03)
    top = Inches(1.56)
    widths = [Inches(1.60), Inches(1.80), Inches(2.10), Inches(6.40)]

    # Header row
    headers = ["Layer", "Tool", "Version / Note", "Role in TDGen Temporal"]
    hx = MARGIN_L
    for w, h in zip(widths, headers):
        rect(slide, hx, top, w, row_h * 0.65, CGI_RED)
        txt(
            slide,
            h,
            hx + Inches(0.06),
            top + Inches(0.05),
            w - Inches(0.10),
            row_h * 0.55,
            size=12,
            bold=True,
            color=CGI_WHITE,
        )
        hx += w + Inches(0.02)

    top += row_h * 0.65 + gap

    CHANGED = {"pandera"}

    for i, (layer, tool, note, role) in enumerate(stack):
        bg = CGI_LIGHT_GREY if i % 2 == 0 else CGI_WHITE
        rx = MARGIN_L
        for j, (w, val) in enumerate(zip(widths, [layer, tool, note, role])):
            rect(slide, rx, top, w, row_h, bg)
            accent = CGI_RED if j == 0 else None
            if j == 0:
                rect(slide, rx, top, Inches(0.04), row_h, CGI_RED)
            color = CGI_RED if (j == 2 and "replaced" in val) else CGI_BODY
            bold = j == 1
            txt(
                slide,
                val,
                rx + Inches(0.08),
                top + Inches(0.07),
                w - Inches(0.12),
                row_h - Inches(0.10),
                size=12,
                bold=bold,
                color=color,
            )
            rx += w + Inches(0.02)
        top += row_h + gap

    add_footer(slide, 6)


# ---------------------------------------------------------------------------
# SLIDE 6  -  Output Model
# ---------------------------------------------------------------------------


def slide_06_output(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Output Model",
        "Every daily run produces INSERT deltas, UPDATE deltas, and an updated state database",
    )

    # Left: file tree
    LEFT_W = Inches(5.20)
    rect(slide, MARGIN_L, Inches(1.56), LEFT_W, Inches(5.10), CGI_LIGHT_GREY)
    rect(slide, MARGIN_L, Inches(1.56), Inches(0.05), Inches(5.10), CGI_RED)

    txt(
        slide,
        "Output Directory Structure",
        MARGIN_L + Inches(0.14),
        Inches(1.62),
        LEFT_W - Inches(0.20),
        Inches(0.32),
        size=13,
        bold=True,
        color=CGI_RED,
    )

    tree = [
        ("output/", False, 0),
        ("  state.db", True, 1),
        ("  deltas/", False, 0),
        ("    2026-04-15/", False, 1),
        ("      inserts/", False, 1),
        ("        TRANSACTION.csv  (580 rows)", True, 2),
        ("        TRANSACTION.json", True, 2),
        ("        AUTHORIZATION.csv  (638 rows)", True, 2),
        ("        STATEMENT.csv  (17 rows)", True, 2),
        ("        DISPUTE.csv  (new)", True, 2),
        ("      updates/", False, 1),
        ("        ACCOUNT.csv  (500 rows)", True, 2),
        ("        DISPUTE.csv  (progressed)", True, 2),
        ("    2026-04-16/  ...", False, 1),
        ("  snapshots/  (optional)", False, 0),
    ]

    ty = Inches(2.02)
    for label, is_file, indent in tree:
        color = CGI_BODY if is_file else CGI_BLACK
        bold = not is_file
        txt(
            slide,
            label,
            MARGIN_L + Inches(0.18),
            ty,
            LEFT_W - Inches(0.26),
            Inches(0.26),
            size=11,
            bold=bold,
            color=color,
        )
        ty += Inches(0.29)

    # Right: three output type cards
    RX = MARGIN_L + LEFT_W + Inches(0.35)
    RW = CONTENT_W - LEFT_W - Inches(0.35)
    CARD_H = Inches(1.52)
    cy = Inches(1.56)

    cards = [
        (
            "INSERT Delta",
            "New rows added today",
            [
                "TRANSACTION  ~580 rows/day",
                "AUTHORIZATION  ~638 rows/day",
                "STATEMENT  ~17 rows/day  (cycle_day match)",
                "DISPUTE, FRAUD_ALERT  (probabilistic)",
            ],
        ),
        (
            "UPDATE Delta",
            "Changed rows — full row snapshot",
            [
                "ACCOUNT  500 rows/day  (balance, status, dates)",
                "DISPUTE  (status progression)",
                "FRAUD_ALERT  (confirmed / false_positive)",
                "CHARGEBACK  (stage advancement)",
            ],
        ),
        (
            "State DB",
            "output/state.db  —  authoritative source of truth",
            [
                "All 15 entity tables  +  16 REF tables",
                "6 temporal state tables  (simulation-internal)",
                "pk_sequences  +  run_log  +  simulation_meta",
                "Queryable with any SQLite tool",
            ],
        ),
    ]

    colors = [CGI_RED, CGI_BLUE, RGBColor(0x1A, 0x7A, 0x3C)]

    for (title, subtitle, bullets), accent in zip(cards, colors):
        rect(slide, RX, cy, RW, CARD_H, CGI_LIGHT_GREY)
        rect(slide, RX, cy, Inches(0.05), CARD_H, accent)
        txt(
            slide,
            title,
            RX + Inches(0.14),
            cy + Inches(0.08),
            RW - Inches(0.20),
            Inches(0.30),
            size=13,
            bold=True,
            color=accent,
        )
        txt(
            slide,
            subtitle,
            RX + Inches(0.14),
            cy + Inches(0.40),
            RW - Inches(0.20),
            Inches(0.22),
            size=11,
            color=CGI_GREY,
        )
        bx = cy + Inches(0.64)
        for b in bullets:
            txt(
                slide,
                f"  {b}",
                RX + Inches(0.14),
                bx,
                RW - Inches(0.20),
                Inches(0.22),
                size=11,
                color=CGI_BODY,
            )
            bx += Inches(0.22)
        cy += CARD_H + Inches(0.27)

    add_footer(slide, 7)


# ---------------------------------------------------------------------------
# SLIDE 7  -  Sample Data: ACCOUNT + CUSTOMER
# ---------------------------------------------------------------------------


def slide_07_samples_1(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Sample Data  -  Account & Customer",
        "Day 7 state  |  Balances updated daily  |  Delinquency tracks missed payments",
    )

    def mini_table(title, cols, rows, x, y, w, col_widths, accent=CGI_RED):
        rect(slide, x, y, w, Inches(0.34), accent)
        txt(
            slide,
            title,
            x + Inches(0.10),
            y + Inches(0.04),
            w - Inches(0.14),
            Inches(0.26),
            size=12,
            bold=True,
            color=CGI_WHITE,
        )
        y += Inches(0.34)
        ROW_H = Inches(0.295)
        # Header
        hx = x
        for cw, ch in zip(col_widths, cols):
            rect(slide, hx, y, cw, ROW_H, RGBColor(0x22, 0x22, 0x22))
            txt(
                slide,
                ch,
                hx + Inches(0.04),
                y + Inches(0.04),
                cw - Inches(0.06),
                ROW_H - Inches(0.06),
                size=10,
                bold=True,
                color=CGI_WHITE,
            )
            hx += cw
        y += ROW_H
        for ri, row in enumerate(rows):
            bg = CGI_LIGHT_GREY if ri % 2 == 0 else CGI_WHITE
            rx = x
            for cw, val in zip(col_widths, row):
                rect(slide, rx, y, cw, ROW_H, bg)
                txt(
                    slide,
                    str(val),
                    rx + Inches(0.04),
                    y + Inches(0.04),
                    cw - Inches(0.06),
                    ROW_H - Inches(0.06),
                    size=10,
                    color=CGI_BODY,
                )
                rx += cw
            y += ROW_H

    # ACCOUNT table
    acc_cols = [
        "account_id",
        "account_number",
        "credit_limit",
        "current_balance",
        "account_status",
        "days_delinquent",
        "risk_score",
    ]
    acc_rows = [
        [1, "6787963526218375", "$11,241", "$6,025", "ACTIVE", 0, 637.5],
        [2, "4642531655329029", "$14,489", "$7,812", "ACTIVE", 0, 573.6],
        [3, "2706517396112117", "$39,303", "$2,473", "DELINQUENT", 45, 714.9],
        [4, "5184729301847562", "$22,000", "$9,841", "ACTIVE", 0, 721.2],
        [5, "3019283746152983", "$8,500", "$3,102", "DELINQUENT", 62, 488.1],
    ]
    acc_widths = [
        Inches(0.75),
        Inches(1.40),
        Inches(0.90),
        Inches(1.10),
        Inches(0.90),
        Inches(0.90),
        Inches(0.75),
    ]
    mini_table(
        "ACCOUNT  (500 rows, updated daily)",
        acc_cols,
        acc_rows,
        MARGIN_L,
        Inches(1.56),
        sum(acc_widths),
        acc_widths,
        CGI_RED,
    )

    # CUSTOMER table
    cust_cols = ["customer_id", "first_name", "last_name", "city", "province", "email", "language"]
    cust_rows = [
        [1, "Benjamin", "Harris", "East Jenny", "PE", "joseph49@example.org", "FR"],
        [2, "Brian", "Simmons", "Port Leahmouth", "NS", "amydickson@example.com", "EN"],
        [3, "Mark", "Davis", "North Jeffberg", "PE", "elainewells@example.org", "FR"],
        [4, "Sarah", "Chen", "Toronto", "ON", "schen@example.ca", "EN"],
        [5, "Michel", "Tremblay", "Montreal", "QC", "mtremblay@example.ca", "FR"],
    ]
    cust_widths = [
        Inches(0.85),
        Inches(0.85),
        Inches(0.90),
        Inches(1.30),
        Inches(0.72),
        Inches(2.05),
        Inches(0.72),
    ]
    mini_table(
        "CUSTOMER  (500 rows — SIN synthetic, no real PII)",
        cust_cols,
        cust_rows,
        MARGIN_L,
        Inches(4.24),
        sum(cust_widths),
        cust_widths,
        CGI_BLUE,
    )

    add_footer(slide, 8)


# ---------------------------------------------------------------------------
# SLIDE 8  -  Sample Data: TRANSACTION + AUTHORIZATION
# ---------------------------------------------------------------------------


def slide_08_samples_2(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Sample Data  -  Transaction & Authorization",
        "~580 transactions generated per day  |  Realistic merchant narratives  |  Approved + declined auths",
    )

    def mini_table(title, cols, rows, x, y, w, col_widths, accent=CGI_RED):
        rect(slide, x, y, w, Inches(0.34), accent)
        txt(
            slide,
            title,
            x + Inches(0.10),
            y + Inches(0.04),
            w - Inches(0.14),
            Inches(0.26),
            size=12,
            bold=True,
            color=CGI_WHITE,
        )
        y += Inches(0.34)
        ROW_H = Inches(0.295)
        hx = x
        for cw, ch in zip(col_widths, cols):
            rect(slide, hx, y, cw, ROW_H, RGBColor(0x22, 0x22, 0x22))
            txt(
                slide,
                ch,
                hx + Inches(0.04),
                y + Inches(0.04),
                cw - Inches(0.06),
                ROW_H - Inches(0.06),
                size=10,
                bold=True,
                color=CGI_WHITE,
            )
            hx += cw
        y += ROW_H
        for ri, row in enumerate(rows):
            bg = CGI_LIGHT_GREY if ri % 2 == 0 else CGI_WHITE
            rx = x
            for cw, val in zip(col_widths, row):
                rect(slide, rx, y, cw, ROW_H, bg)
                txt(
                    slide,
                    str(val),
                    rx + Inches(0.04),
                    y + Inches(0.04),
                    cw - Inches(0.06),
                    ROW_H - Inches(0.06),
                    size=10,
                    color=CGI_BODY,
                )
                rx += cw
            y += ROW_H

    # TRANSACTION
    tx_cols = ["txn_id", "account_id", "date", "amount", "type", "description", "status"]
    tx_rows = [
        [1, 3, "2026-04-15 08:09", "$32.31", "PAYMENT", "COSTCO WHOLESALE #1409", "posted"],
        [2, 3, "2026-04-15 22:04", "$183.09", "PURCHASE", "SHOPIFY* JOHNSON STORE", "posted"],
        [3, 4, "2026-04-15 17:46", "$749.59", "PAYMENT", "APPLE.COM/BILL", "posted"],
        [47, 22, "2026-04-16 10:12", "$12.50", "PURCHASE", "TIM HORTONS #3821", "posted"],
        [119, 88, "2026-04-17 19:33", "$1,240", "PURCHASE", "AMZN MKTP CA*847291", "posted"],
    ]
    tx_widths = [
        Inches(0.60),
        Inches(0.80),
        Inches(1.30),
        Inches(0.80),
        Inches(0.90),
        Inches(2.20),
        Inches(0.80),
    ]
    mini_table(
        "TRANSACTION  (~580 new rows/day)",
        tx_cols,
        tx_rows,
        MARGIN_L,
        Inches(1.56),
        sum(tx_widths),
        tx_widths,
        CGI_RED,
    )

    # AUTHORIZATION
    auth_cols = [
        "auth_id",
        "account_id",
        "amount",
        "response",
        "approval_code",
        "network",
        "channel",
        "decline_reason",
    ]
    auth_rows = [
        [1, 3, "$32.31", "00 Approved", "559407", "AmexNet", "in_store", "-"],
        [2, 3, "$183.09", "00 Approved", "316475", "BankNet", "atm", "-"],
        [3, 4, "$749.59", "00 Approved", "835030", "BankNet", "online", "-"],
        [18, 44, "$520.00", "51 Declined", "-", "VisaNet", "in_store", "insufficient_funds"],
        [31, 112, "$89.99", "54 Declined", "-", "BankNet", "online", "expired_card"],
    ]
    auth_widths = [
        Inches(0.60),
        Inches(0.80),
        Inches(0.80),
        Inches(1.10),
        Inches(1.00),
        Inches(0.90),
        Inches(0.80),
        Inches(1.60),
    ]
    mini_table(
        "AUTHORIZATION  (~638 new rows/day  |  includes declined)",
        auth_cols,
        auth_rows,
        MARGIN_L,
        Inches(4.24),
        sum(auth_widths),
        auth_widths,
        RGBColor(0x1A, 0x7A, 0x3C),
    )

    add_footer(slide, 9)


# ---------------------------------------------------------------------------
# SLIDE 9  -  Sample Data: Dispute + Fraud + Collection
# ---------------------------------------------------------------------------


def slide_09_samples_3(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Sample Data  -  Dispute, Fraud Alert & Collection Case",
        "Lifecycle entities progress through states each day  |  7-day sample shown",
    )

    def mini_table(title, cols, rows, x, y, w, col_widths, accent=CGI_RED):
        rect(slide, x, y, w, Inches(0.34), accent)
        txt(
            slide,
            title,
            x + Inches(0.10),
            y + Inches(0.04),
            w - Inches(0.14),
            Inches(0.26),
            size=12,
            bold=True,
            color=CGI_WHITE,
        )
        y += Inches(0.34)
        ROW_H = Inches(0.275)
        hx = x
        for cw, ch in zip(col_widths, cols):
            rect(slide, hx, y, cw, ROW_H, RGBColor(0x22, 0x22, 0x22))
            txt(
                slide,
                ch,
                hx + Inches(0.04),
                y + Inches(0.04),
                cw - Inches(0.06),
                ROW_H - Inches(0.06),
                size=10,
                bold=True,
                color=CGI_WHITE,
            )
            hx += cw
        y += ROW_H
        for ri, row in enumerate(rows):
            bg = CGI_LIGHT_GREY if ri % 2 == 0 else CGI_WHITE
            rx = x
            for cw, val in zip(col_widths, row):
                rect(slide, rx, y, cw, ROW_H, bg)
                txt(
                    slide,
                    str(val),
                    rx + Inches(0.04),
                    y + Inches(0.03),
                    cw - Inches(0.06),
                    ROW_H - Inches(0.05),
                    size=10,
                    color=CGI_BODY,
                )
                rx += cw
            y += ROW_H

    # DISPUTE
    d_cols = ["dispute_id", "txn_id", "type", "amount", "opened", "status", "resolution"]
    d_rows = [
        [1, 43, "DUPLICATE", "$37.75", "2026-04-15", "RESOLVED", "APPROVED"],
        [2, 541, "QUALITY", "$170.78", "2026-04-15", "RESOLVED", "DENIED"],
        [3, 696, "WRONG_AMOUNT", "$724.13", "2026-04-16", "INVESTIGATING", "-"],
        [4, 812, "FRAUD", "$489.00", "2026-04-17", "INVESTIGATING", "-"],
        [5, 1024, "NOT_RECEIVED", "$62.50", "2026-04-18", "OPEN", "-"],
    ]
    d_widths = [
        Inches(0.80),
        Inches(0.65),
        Inches(1.30),
        Inches(0.80),
        Inches(1.00),
        Inches(1.20),
        Inches(0.90),
    ]
    mini_table(
        "DISPUTE  (8 total  |  progressing daily)",
        d_cols,
        d_rows,
        MARGIN_L,
        Inches(1.56),
        sum(d_widths),
        d_widths,
        CGI_RED,
    )

    # FRAUD_ALERT
    fa_cols = [
        "alert_id",
        "account_id",
        "source",
        "type",
        "risk_score",
        "status",
        "action",
        "resolved",
    ]
    fa_rows = [
        [1, 22, "3DS", "LARGE_TXN", 621, "CONFIRMED", "block_card", "2026-04-17"],
        [2, 29, "CardGuard", "LARGE_TXN", 226, "CONFIRMED", "block_card", "2026-04-20"],
    ]
    fa_widths = [
        Inches(0.65),
        Inches(0.80),
        Inches(0.85),
        Inches(1.10),
        Inches(0.80),
        Inches(1.00),
        Inches(0.90),
        Inches(0.95),
    ]
    mini_table(
        "FRAUD_ALERT  (2 total  |  both confirmed, cards blocked)",
        fa_cols,
        fa_rows,
        MARGIN_L,
        Inches(3.84),
        sum(fa_widths),
        fa_widths,
        RGBColor(0xCC, 0x44, 0x00),
    )

    # COLLECTION_CASE
    cc_cols = [
        "case_id",
        "account_id",
        "opened",
        "bucket",
        "past_due",
        "total_owed",
        "status",
        "next_action",
    ]
    cc_rows = [
        [1, 3, "2026-04-15", "B2", "$49.46", "$2,473", "ACTIVE", "initial_contact"],
        [2, 42, "2026-04-15", "B2", "$696.89", "$34,844", "ACTIVE", "initial_contact"],
        [3, 81, "2026-04-15", "B3", "$417.80", "$20,890", "ACTIVE", "initial_contact"],
        [27, 498, "2026-04-21", "B1", "$220.50", "$11,025", "ACTIVE", "initial_contact"],
    ]
    cc_widths = [
        Inches(0.60),
        Inches(0.80),
        Inches(0.95),
        Inches(0.60),
        Inches(0.80),
        Inches(0.80),
        Inches(0.80),
        Inches(1.45),
    ]
    mini_table(
        "COLLECTION_CASE  (27 total  |  auto-opened when days_delinquent >= 30)",
        cc_cols,
        cc_rows,
        MARGIN_L,
        Inches(5.32),
        sum(cc_widths),
        cc_widths,
        CGI_BLUE,
    )

    add_footer(slide, 10)


# ---------------------------------------------------------------------------
# SLIDE 10  -  7-Day Run Summary
# ---------------------------------------------------------------------------


def slide_10_results(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "7-Day Run Summary",
        "init (Day 0)  +  advance x7  |  ~0.5 seconds per day  |  500 accounts",
    )

    # Left: per-day table
    LEFT_W = Inches(5.80)
    rect(slide, MARGIN_L, Inches(1.56), LEFT_W, Inches(4.80), CGI_LIGHT_GREY)
    rect(slide, MARGIN_L, Inches(1.56), Inches(0.05), Inches(4.80), CGI_RED)

    txt(
        slide,
        "Daily Run Log",
        MARGIN_L + Inches(0.14),
        Inches(1.62),
        LEFT_W - Inches(0.20),
        Inches(0.30),
        size=13,
        bold=True,
        color=CGI_RED,
    )

    day_cols = ["Date", "Inserts", "Txns", "Stmts", "Updates", "Dur"]
    day_widths = [
        Inches(1.20),
        Inches(0.80),
        Inches(0.70),
        Inches(0.60),
        Inches(0.90),
        Inches(0.50),
    ]
    day_data = [
        ["2026-04-14", "1,733", "-", "-", "0", "0.7s"],
        ["2026-04-15", "1,268", "580", "17", "961", "0.4s"],
        ["2026-04-16", "1,239", "580", "16", "991", "0.5s"],
        ["2026-04-17", "1,223", "571", "16", "992", "0.4s"],
        ["2026-04-18", "1,213", "572", "15", "991", "0.5s"],
        ["2026-04-19", "1,189", "559", "16", "992", "0.5s"],
        ["2026-04-20", "1,241", "594", "15", "994", "0.5s"],
        ["2026-04-21", "1,209", "573", "24", "992", "0.4s"],
    ]

    ROW_H = Inches(0.35)
    ty = Inches(2.00)
    hx = MARGIN_L + Inches(0.08)
    for cw, ch in zip(day_widths, day_cols):
        rect(slide, hx, ty, cw, ROW_H, RGBColor(0x22, 0x22, 0x22))
        txt(
            slide,
            ch,
            hx + Inches(0.04),
            ty + Inches(0.05),
            cw - Inches(0.06),
            ROW_H - Inches(0.08),
            size=11,
            bold=True,
            color=CGI_WHITE,
        )
        hx += cw
    ty += ROW_H

    for i, row in enumerate(day_data):
        bg = CGI_LIGHT_GREY if i % 2 == 0 else CGI_WHITE
        if i == 0:
            bg = RGBColor(0xFA, 0xE8, 0xE8)  # init row highlight
        rx = MARGIN_L + Inches(0.08)
        for cw, val in zip(day_widths, row):
            rect(slide, rx, ty, cw, ROW_H, bg)
            txt(
                slide,
                val,
                rx + Inches(0.04),
                ty + Inches(0.05),
                cw - Inches(0.06),
                ROW_H - Inches(0.08),
                size=11,
                color=CGI_BODY,
            )
            rx += cw
        ty += ROW_H

    # Right: cumulative stats
    RX = MARGIN_L + LEFT_W + Inches(0.40)
    RW = CONTENT_W - LEFT_W - Inches(0.40)

    txt(
        slide,
        "Cumulative State (Day 7)",
        RX,
        Inches(1.62),
        RW,
        Inches(0.30),
        size=13,
        bold=True,
        color=CGI_RED,
    )

    stats = [
        ("AUTHORIZATION", "4,397"),
        ("TRANSACTION", "4,029"),
        ("ACCOUNT", "500"),
        ("CUSTOMER", "500"),
        ("CARD", "500"),
        ("MERCHANT", "200"),
        ("STATEMENT", "119"),
        ("COLLECTION_CASE", "27"),
        ("PRODUCT_DEFINITION", "20"),
        ("PROVIDER", "10"),
        ("DISPUTE", "8"),
        ("CLIENT", "3"),
        ("FRAUD_ALERT", "2"),
    ]

    sy = Inches(2.04)
    for i, (tbl, count) in enumerate(stats):
        bg = CGI_LIGHT_GREY if i % 2 == 0 else CGI_WHITE
        rect(slide, RX, sy, RW, Inches(0.30), bg)
        rect(slide, RX, sy, Inches(0.04), Inches(0.30), CGI_RED)
        txt(
            slide,
            tbl,
            RX + Inches(0.10),
            sy + Inches(0.04),
            RW * 0.65,
            Inches(0.22),
            size=11,
            color=CGI_BODY,
        )
        txt(
            slide,
            count,
            RX + RW * 0.72,
            sy + Inches(0.04),
            RW * 0.26,
            Inches(0.22),
            size=11,
            bold=True,
            color=CGI_RED,
            align=PP_ALIGN.RIGHT,
        )
        sy += Inches(0.30)

    # Bottom total bar
    BAR_Y = Inches(6.52)
    rect(slide, MARGIN_L, BAR_Y, CONTENT_W, Inches(0.44), CGI_RED)
    txt(
        slide,
        "8 runs  (init + 7 days)  |  10,207 total inserts  |  6,921 total updates  |  ~0.5s per day  |  output/state.db",
        MARGIN_L + Inches(0.20),
        BAR_Y + Inches(0.06),
        CONTENT_W - Inches(0.40),
        Inches(0.34),
        size=13,
        bold=True,
        color=CGI_WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, 11)


# ---------------------------------------------------------------------------
# SLIDE 11  -  CLI Reference
# ---------------------------------------------------------------------------


def slide_11_cli(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "CLI Reference",
        "All modes run from the terminal — no GUI, no configuration files to edit for daily use",
    )

    commands = [
        (
            "init",
            ".venv/Scripts/python -m tdgen_temporal.cli init",
            "Seed the Day 0 population. Run once. Creates output/state.db with 500 accounts,\n"
            "500 customers, 500 cards, 200 merchants, and all reference tables.",
        ),
        (
            "advance",
            ".venv/Scripts/python -m tdgen_temporal.cli advance --days 30",
            "Advance the simulation by N days from the current state. Writes delta files to\n"
            "output/deltas/YYYY-MM-DD/inserts/ and updates/ each day.",
        ),
        (
            "backfill",
            ".venv/Scripts/python -m tdgen_temporal.cli backfill --from 2026-01-01 --to 2026-04-14",
            "Generate a full date range of history in sequence. Useful for creating a rich\n"
            "starting dataset before handing off to QA.",
        ),
        (
            "status",
            ".venv/Scripts/python -m tdgen_temporal.cli status",
            "Show current simulation date, total runs, and row counts for all entity tables.\n"
            "No side effects — safe to run at any time.",
        ),
    ]

    BOX_H = Inches(1.36)
    GAP = Inches(0.14)
    top = Inches(1.56)

    for cmd, command, desc in commands:
        rect(slide, MARGIN_L, top, CONTENT_W, BOX_H, CGI_LIGHT_GREY)
        rect(slide, MARGIN_L, top, Inches(0.06), BOX_H, CGI_RED)

        txt(
            slide,
            cmd,
            MARGIN_L + Inches(0.14),
            top + Inches(0.08),
            Inches(0.80),
            Inches(0.28),
            size=12,
            bold=True,
            color=CGI_RED,
        )

        # Command chip
        rect(
            slide,
            MARGIN_L + Inches(1.02),
            top + Inches(0.06),
            Inches(7.60),
            Inches(0.36),
            RGBColor(0x22, 0x22, 0x22),
        )
        txt(
            slide,
            command,
            MARGIN_L + Inches(1.10),
            top + Inches(0.08),
            Inches(7.44),
            Inches(0.30),
            size=12,
            bold=True,
            color=RGBColor(0xFF, 0xCC, 0x44),
        )

        txt(
            slide,
            desc,
            MARGIN_L + Inches(0.14),
            top + Inches(0.52),
            CONTENT_W - Inches(0.24),
            Inches(0.76),
            size=13,
            color=CGI_BODY,
        )

        top += BOX_H + GAP

    add_footer(slide, 12)


# ---------------------------------------------------------------------------
# SLIDE 12  -  Next Steps
# ---------------------------------------------------------------------------


def slide_12_nextsteps(prs):
    slide = blank_slide(prs)
    slide_header(
        slide, "Next Steps", "TDGen Temporal  |  Phase 1 complete  |  TD Bank / CGI engagement"
    )

    done = [
        "Day-over-day temporal progression — accounts, balances, delinquency",
        "6 entity state machines — ACCOUNT, CARD, DISPUTE, FRAUD_ALERT, CHARGEBACK, COLLECTION_CASE",
        "Daily transaction + authorization generation (~580 txns/day per 500 accounts)",
        "Statement generation on account cycle_day",
        "Score record monthly refresh with drift based on account health",
        "INSERT and UPDATE delta files per day  (CSV + JSON)",
        "SQLite state store with WAL mode — persistent, queryable, restartable",
        "CLI: init / advance / backfill / status",
        "config/scenario.yaml — all rates configurable without code changes",
    ]

    next_items = [
        "Validate with TD QA team — confirm rates match expected fraud lifecycle volumes",
        "New account opens — daily probabilistic seeding of new ACCOUNT + CUSTOMER + CARD",
        "MCP server wrapper — expose init/advance/status as MCP tools for Claude Code",
        "GitHub Actions workflow — scheduled nightly advance + delta archival",
        "Scenario presets — high-fraud day, payment cycle stress, chargeoff wave",
        "PostgreSQL provisioner — SQLAlchemy bulk insert for staging environment",
    ]

    LW = Inches(5.90)
    RX = MARGIN_L + LW + Inches(0.43)
    RW = CONTENT_W - LW - Inches(0.43)

    # Done column
    rect(slide, MARGIN_L, Inches(1.56), LW, Inches(0.36), RGBColor(0x1A, 0x7A, 0x3C))
    txt(
        slide,
        "Phase 1 Complete",
        MARGIN_L + Inches(0.10),
        Inches(1.60),
        LW - Inches(0.14),
        Inches(0.28),
        size=13,
        bold=True,
        color=CGI_WHITE,
    )

    dy = Inches(2.00)
    for item in done:
        rect(slide, MARGIN_L, dy, LW, Inches(0.44), CGI_LIGHT_GREY)
        rect(slide, MARGIN_L, dy, Inches(0.04), Inches(0.44), RGBColor(0x1A, 0x7A, 0x3C))
        txt(
            slide,
            "OK  " + item,
            MARGIN_L + Inches(0.12),
            dy + Inches(0.06),
            LW - Inches(0.18),
            Inches(0.34),
            size=11,
            color=CGI_BODY,
        )
        dy += Inches(0.46)

    # Next column
    rect(slide, RX, Inches(1.56), RW, Inches(0.36), CGI_RED)
    txt(
        slide,
        "Phase 2 — Planned",
        RX + Inches(0.10),
        Inches(1.60),
        RW - Inches(0.14),
        Inches(0.28),
        size=13,
        bold=True,
        color=CGI_WHITE,
    )

    ny = Inches(2.00)
    for item in next_items:
        rect(slide, RX, ny, RW, Inches(0.52), CGI_LIGHT_GREY)
        rect(slide, RX, ny, Inches(0.04), Inches(0.52), CGI_RED)
        txt(
            slide,
            item,
            RX + Inches(0.12),
            ny + Inches(0.08),
            RW - Inches(0.18),
            Inches(0.40),
            size=11,
            color=CGI_BODY,
        )
        ny += Inches(0.54)

    add_footer(slide, 13)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


def build_presentation(output_path: str = "docs/TDGen_Temporal_v1.pptx"):
    prs = new_prs()
    slide_01_cover(prs)
    slide_02_what(prs)
    slide_03_scope(prs)
    slide_03_how(prs)
    slide_04_state_machines(prs)
    slide_05_stack(prs)
    slide_06_output(prs)
    slide_07_samples_1(prs)
    slide_08_samples_2(prs)
    slide_09_samples_3(prs)
    slide_10_results(prs)
    slide_11_cli(prs)
    slide_12_nextsteps(prs)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(prs.slides)} slides)")


# ===========================================================================
# V4  —  30-Day Run  (init 2026-03-31  +  advance x30  →  2026-04-30)
# ===========================================================================

_V4_RUN_LOG = [
    # (date,           inserts, txns, updates, dur)
    ("2026-03-31 init", 1733, 0, 0, "0.4s"),
    ("2026-04-01", 1752, 574, 961, "0.4s"),
    ("2026-04-02", 1205, 569, 989, "0.3s"),
    ("2026-04-03", 1208, 565, 988, "0.3s"),
    ("2026-04-04", 1215, 568, 991, "0.4s"),
    ("2026-04-05", 1190, 562, 988, "0.4s"),
    ("2026-04-06", 1256, 595, 986, "0.3s"),
    ("2026-04-07", 1228, 584, 986, "0.3s"),
    ("2026-04-08", 1249, 586, 987, "0.3s"),
    ("2026-04-09", 1199, 569, 987, "0.4s"),
    ("2026-04-10", 1169, 555, 989, "0.3s"),
    ("2026-04-11", 1240, 580, 988, "0.3s"),
    ("2026-04-12", 1317, 618, 993, "0.4s"),
    ("2026-04-13", 1272, 592, 993, "0.3s"),
    ("2026-04-14", 1264, 598, 993, "0.4s"),
    ("2026-04-15", 1196, 562, 994, "0.3s"),
    ("2026-04-16", 1240, 585, 994, "0.4s"),
    ("2026-04-17", 1241, 580, 997, "0.4s"),
    ("2026-04-18", 1255, 598, 999, "0.4s"),
    ("2026-04-19", 1260, 594, 996, "0.4s"),
    ("2026-04-20", 1154, 545, 996, "0.4s"),
    ("2026-04-21", 1319, 622, 999, "0.4s"),
    ("2026-04-22", 1223, 576, 999, "0.4s"),
    ("2026-04-23", 1188, 565, 995, "0.4s"),
    ("2026-04-24", 1195, 552, 991, "0.3s"),
    ("2026-04-25", 1273, 601, 992, "0.3s"),
    ("2026-04-26", 1285, 607, 988, "0.4s"),
    ("2026-04-27", 1228, 576, 994, "0.4s"),
    ("2026-04-28", 1187, 553, 994, "0.4s"),
    ("2026-04-29", 1289, 621, 995, "0.4s"),
    ("2026-04-30", 1191, 568, 994, "0.4s"),
]

_V4_FINAL_STATE = [
    ("AUTHORIZATION", 19003),
    ("TRANSACTION", 17420),
    ("ACCOUNT", 500),
    ("CUSTOMER", 500),
    ("CARD", 500),
    ("STATEMENT", 494),
    ("SCORE_RECORD", 494),
    ("MERCHANT", 200),
    ("COLLECTION_CASE", 27),
    ("DISPUTE", 35),
    ("PRODUCT_DEFINITION", 20),
    ("PROVIDER", 10),
    ("FRAUD_ALERT", 11),
    ("CHARGEBACK", 4),
    ("CLIENT", 3),
]


# ---------------------------------------------------------------------------
# V4 helper: mini_table
# ---------------------------------------------------------------------------


def _v4_mini_table(slide, title, cols, rows, x, y, col_widths, accent=None):
    """Draw a titled mini-table; returns bottom y coordinate."""
    if accent is None:
        accent = CGI_RED
    w = sum(col_widths)
    rect(slide, x, y, w, Inches(0.32), accent)
    txt(
        slide,
        title,
        x + Inches(0.09),
        y + Inches(0.04),
        w - Inches(0.14),
        Inches(0.24),
        size=11,
        bold=True,
        color=CGI_WHITE,
    )
    y += Inches(0.32)
    ROW_H = Inches(0.265)
    hx = x
    for cw, ch in zip(col_widths, cols):
        rect(slide, hx, y, cw, ROW_H, RGBColor(0x22, 0x22, 0x22))
        txt(
            slide,
            ch,
            hx + Inches(0.04),
            y + Inches(0.03),
            cw - Inches(0.06),
            ROW_H - Inches(0.05),
            size=9,
            bold=True,
            color=CGI_WHITE,
        )
        hx += cw
    y += ROW_H
    for ri, row in enumerate(rows):
        bg = CGI_LIGHT_GREY if ri % 2 == 0 else CGI_WHITE
        rx = x
        for cw, val in zip(col_widths, row):
            rect(slide, rx, y, cw, ROW_H, bg)
            txt(
                slide,
                str(val),
                rx + Inches(0.04),
                y + Inches(0.03),
                cw - Inches(0.06),
                ROW_H - Inches(0.05),
                size=9,
                color=CGI_BODY,
            )
            rx += cw
        y += ROW_H
    return y


# ---------------------------------------------------------------------------
# V4 SLIDE 1  -  Cover
# ---------------------------------------------------------------------------


def slide_v4_01_cover(prs):
    slide = blank_slide(prs)
    cornerstone(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, CGI_RED)

    txt(
        slide,
        "Test Data Generation",
        Inches(0.55),
        Inches(1.30),
        Inches(6.0),
        Inches(0.72),
        size=38,
        bold=True,
        color=CGI_BLACK,
    )
    txt(
        slide,
        "with Claude Code",
        Inches(0.55),
        Inches(2.05),
        Inches(6.0),
        Inches(0.72),
        size=38,
        bold=True,
        color=CGI_BLACK,
    )
    txt(
        slide,
        "Day-Over-Day Synthetic Test Data Generation",
        Inches(0.55),
        Inches(2.84),
        Inches(6.0),
        Inches(0.48),
        size=18,
        bold=False,
        color=CGI_BODY,
    )

    rect(slide, Inches(0.55), Inches(3.42), Inches(5.80), Inches(0.04), CGI_RED)

    multi_para(
        slide,
        [
            ("TSYS TS2  |  Credit Card Fraud Lifecycle", 15, False, CGI_BODY),
            ("TD Bank  |  CGI Consulting Engagement", 15, False, CGI_BODY),
            ("Frederick Ferguson, CGI DCS", 13, False, CGI_GREY),
            ("April 2026", 13, False, CGI_GREY),
        ],
        Inches(0.55),
        Inches(3.58),
        Inches(6.0),
        Inches(1.20),
        space_pt=5,
    )

    add_footer(slide, 1)


# ---------------------------------------------------------------------------
# V4 SLIDE 2  -  Overview  (reused from v2)
# ---------------------------------------------------------------------------
# slide_02_what(prs) is called directly


# ---------------------------------------------------------------------------
# V4 SLIDE 3  -  Scope: TSYS TS2 (text + stat tiles)
# ---------------------------------------------------------------------------


def slide_v4_03_scope_text(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Scope  \u2014  TSYS TS2 Credit Card Processing Schema",
        "Why TS2?  Directly relevant to TD Bank fraud management  |  Enterprise-scale stress-test for TDGen",
    )

    rect(slide, MARGIN_L, Inches(1.56), CONTENT_W, Inches(2.24), CGI_LIGHT_GREY)
    rect(slide, MARGIN_L, Inches(1.56), Inches(0.06), Inches(2.24), CGI_RED)

    multi_para(
        slide,
        [
            ("Why TSYS TS2?", 13, True, CGI_BLACK),
            (
                "We selected TSYS TS2 Credit Card Processing as our test schema because it is directly relevant to "
                "TD Bank\u2019s fraud management workflow and provides an enterprise-scale stress-test for TDGen.",
                12,
                False,
                CGI_BODY,
            ),
            ("", 5, False, CGI_BODY),
            ("Schema Approach", 13, True, CGI_BLACK),
            (
                "Since TS2\u2019s schema is proprietary, we constructed a proxy by synthesizing available sources from "
                "TSYS documentation, Mphasis integration guides, FICO risk model references, and ISO 8583 standards. "
                "The resulting approximation is used exclusively for this POC and contains:",
                12,
                False,
                CGI_BODY,
            ),
        ],
        MARGIN_L + Inches(0.14),
        Inches(1.68),
        CONTENT_W - Inches(0.24),
        Inches(2.04),
        space_pt=5,
    )

    stats = [
        ("31", "Tables"),
        ("15", "Entity\nTables"),
        ("16", "Lookup\nTables"),
        ("285", "Columns"),
        ("47", "Relationships"),
    ]
    TILE_GAP = Inches(0.12)
    TILE_W = (CONTENT_W - 4 * TILE_GAP) // 5
    TILE_H = Inches(1.72)
    tx = MARGIN_L
    ty = Inches(3.96)

    for num, label in stats:
        rect(slide, tx, ty, TILE_W, TILE_H, CGI_BLACK)
        txt(
            slide,
            num,
            tx,
            ty + Inches(0.18),
            TILE_W,
            Inches(0.90),
            size=34,
            bold=True,
            color=CGI_WHITE,
            align=PP_ALIGN.CENTER,
        )
        txt(
            slide,
            label,
            tx,
            ty + Inches(1.12),
            TILE_W,
            Inches(0.50),
            size=10,
            color=RGBColor(0xBB, 0xBB, 0xBB),
            align=PP_ALIGN.CENTER,
        )
        tx += TILE_W + TILE_GAP

    add_footer(slide, 3)


# ---------------------------------------------------------------------------
# V4 SLIDE 4  -  Scope: Six State Machines (simple list layout)
# ---------------------------------------------------------------------------


def slide_v4_04_state_machines(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Scope  \u2014  Entity State Machines",
        "Six entities advance through defined lifecycle states each simulation day  |  Rates configurable in scenario.yaml",
    )

    machines = [
        (
            "ACCOUNT",
            "ACTIVE  \u2192  DELINQUENT  \u2192  CHARGEOFF  \u2192  CLOSED",
            "Missed payments age delinquency; payment received restores ACTIVE; 180+ days triggers CHARGEOFF",
            RGBColor(0x00, 0x5B, 0x99),
        ),
        (
            "CARD",
            "ACTIVE  \u2192  BLOCKED  \u2192  EXPIRED  \u2192  CLOSED",
            "Confirmed fraud alert blocks card; false-positive restores; expiry triggers replacement",
            RGBColor(0x00, 0x6E, 0x51),
        ),
        (
            "DISPUTE",
            "OPEN  \u2192  INVESTIGATING  \u2192  RESOLVED  \u2192  CLOSED",
            "Auto-progressed each day; customer may WITHDRAW within 3 days; resolution in 7+ days",
            RGBColor(0x1A, 0x7A, 0x3C),
        ),
        (
            "FRAUD ALERT",
            "OPEN  \u2192  UNDER_REVIEW  \u2192  CONFIRMED  |  FALSE_POSITIVE",
            "95% next-day pickup; CONFIRMED rate 30% (blocks card); FALSE_POSITIVE 70% (restores card)",
            CGI_ORANGE,
        ),
        (
            "CHARGEBACK",
            "FIRST_CHARGEBACK  \u2192  REPRESENTMENT  \u2192  WON  |  LOST",
            "Escalated from RESOLVED disputes; merchant representment window 10 days",
            RGBColor(0x5E, 0x1A, 0x6E),
        ),
        (
            "COLLECTION",
            "B1  \u2192  B2  \u2192  B3  \u2192  B4  \u2192  CHARGEOFF",
            "Auto-opened at 30 days delinquent; buckets age every 30 days; CHARGEOFF at 120+",
            RGBColor(0x1A, 0x3C, 0x5E),
        ),
    ]

    ROW_H = Inches(0.72)
    ROW_G = Inches(0.04)
    NAME_W = Inches(1.50)
    FLOW_W = Inches(4.60)
    DESC_W = CONTENT_W - NAME_W - FLOW_W
    ty = Inches(1.56)

    for name, flow, desc, color in machines:
        rect(slide, MARGIN_L, ty, CONTENT_W, ROW_H, CGI_LIGHT_GREY)
        rect(slide, MARGIN_L, ty, Inches(0.06), ROW_H, color)
        # Name
        txt(
            slide,
            name,
            MARGIN_L + Inches(0.12),
            ty + Inches(0.22),
            NAME_W - Inches(0.12),
            Inches(0.30),
            size=13,
            bold=True,
            color=color,
        )
        # State flow
        txt(
            slide,
            flow,
            MARGIN_L + NAME_W,
            ty + Inches(0.06),
            FLOW_W - Inches(0.10),
            Inches(0.28),
            size=11,
            bold=True,
            color=CGI_BLACK,
        )
        # Description
        txt(
            slide,
            desc,
            MARGIN_L + NAME_W,
            ty + Inches(0.36),
            FLOW_W + DESC_W - Inches(0.14),
            Inches(0.28),
            size=10,
            color=CGI_GREY,
        )
        ty += ROW_H + ROW_G

    add_footer(slide, 4)


# ---------------------------------------------------------------------------
# V4 SLIDE 5  -  9-Step Daily Pipeline
# ---------------------------------------------------------------------------


def slide_v4_05_pipeline(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "The TDGen Daily Data Generation Pipeline",
        "Each day, active accounts are loaded and every entity is advanced in dependency order",
    )

    txt(
        slide,
        "Each day, TDGen loads active accounts and advances each entity through its lifecycle \u2014 "
        "updating accounts and cards, generating transactions, statements, and scores, "
        "then progressing disputes, fraud alerts, chargebacks, and collection cases \u2014 "
        "before persisting all changes and exporting delta files.",
        MARGIN_L,
        Inches(1.56),
        CONTENT_W,
        Inches(0.54),
        size=12,
        color=CGI_BODY,
    )

    steps = [
        (
            "1",
            "Update Accounts",
            "Balance updates, delinquency aging, ACTIVE \u2192 DELINQUENT \u2192 CHARGEOFF",
        ),
        ("2", "Update Cards", "Expiry checks, BLOCKED cards, replacement issuance"),
        ("3", "Update Scores", "Monthly risk score drift (PD) based on account health"),
        ("4", "Generate Transactions", "Purchases, payments, fees, cash advances  (~581/day)"),
        (
            "5",
            "Generate Authorizations",
            "One auth per transaction + probabilistic declined-only auths  (~633/day)",
        ),
        ("6", "Generate Statements", "Triggered when run_date.day == account.cycle_day  (~18/day)"),
        (
            "7",
            "Progress Disputes",
            "OPEN \u2192 INVESTIGATING \u2192 RESOLVED \u2192 CLOSED  (35 total, 30-day run)",
        ),
        (
            "8",
            "Progress Fraud Alerts",
            "OPEN \u2192 UNDER_REVIEW \u2192 CONFIRMED | FALSE_POSITIVE  (11 total)",
        ),
        (
            "9",
            "Progress Chargebacks & Collections",
            "CHARGEBACK: stage advancement (4 total)  |  COLLECTION_CASE: bucket progression (27 total)",
        ),
    ]

    col_w = Inches(5.90)
    row_h = Inches(0.415)
    gap = Inches(0.04)
    top = Inches(2.20)
    left1 = MARGIN_L
    left2 = MARGIN_L + col_w + Inches(0.43)

    # Steps 1-8 in two-column layout
    for i, (num, heading, detail) in enumerate(steps[:8]):
        col = i % 2
        row = i // 2
        x = left1 if col == 0 else left2
        y = top + row * (row_h + gap)
        bg = CGI_LIGHT_GREY if row % 2 == 0 else CGI_WHITE

        rect(slide, x, y, col_w, row_h, bg)
        rect(slide, x, y, Inches(0.04), row_h, CGI_RED)
        txt(
            slide,
            num,
            x + Inches(0.10),
            y + Inches(0.06),
            Inches(0.30),
            Inches(0.28),
            size=13,
            bold=True,
            color=CGI_RED,
        )
        txt(
            slide,
            heading,
            x + Inches(0.44),
            y + Inches(0.06),
            Inches(1.75),
            Inches(0.28),
            size=12,
            bold=True,
            color=CGI_BLACK,
        )
        txt(
            slide,
            detail,
            x + Inches(2.24),
            y + Inches(0.06),
            col_w - Inches(2.32),
            Inches(0.28),
            size=11,
            color=CGI_BODY,
        )

    # Step 9  —  full-width row
    y9 = top + 4 * (row_h + gap)
    rect(slide, left1, y9, CONTENT_W, row_h, CGI_LIGHT_GREY)
    rect(slide, left1, y9, Inches(0.04), row_h, CGI_RED)
    txt(
        slide,
        "9",
        left1 + Inches(0.10),
        y9 + Inches(0.06),
        Inches(0.30),
        Inches(0.28),
        size=13,
        bold=True,
        color=CGI_RED,
    )
    txt(
        slide,
        "Progress Chargebacks & Collections",
        left1 + Inches(0.44),
        y9 + Inches(0.06),
        Inches(3.00),
        Inches(0.28),
        size=12,
        bold=True,
        color=CGI_BLACK,
    )
    txt(
        slide,
        "CHARGEBACK: stage advancement (4 total)  |  COLLECTION_CASE: bucket progression (27 total, 30-day run)",
        left1 + Inches(3.60),
        y9 + Inches(0.06),
        CONTENT_W - Inches(3.68),
        Inches(0.28),
        size=11,
        color=CGI_BODY,
    )

    add_footer(slide, 5)


# ---------------------------------------------------------------------------
# V4 SLIDE 6  -  Output Model  +  Account & Customer Sample  (COMBINED)
# ---------------------------------------------------------------------------


def slide_v4_06_output_sample(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Output Model  +  Account & Customer Sample",
        "init 2026-03-31 (Day 0)  +  advance \u00d730  (2026-04-01 \u2192 2026-04-30)  |  500 accounts  |  ~0.35s/day",
    )

    PANEL_TOP = Inches(1.56)
    PANEL_H = Inches(5.10)
    LEFT_W = Inches(4.62)
    GAP = Inches(0.22)
    RIGHT_X = MARGIN_L + LEFT_W + GAP
    RIGHT_W = CONTENT_W - LEFT_W - GAP

    # ── Left: output directory tree ─────────────────────────────────────────
    rect(slide, MARGIN_L, PANEL_TOP, LEFT_W, PANEL_H, CGI_LIGHT_GREY)
    rect(slide, MARGIN_L, PANEL_TOP, Inches(0.06), PANEL_H, CGI_RED)

    txt(
        slide,
        "OUTPUT DIRECTORY  (30-day run)",
        MARGIN_L + Inches(0.14),
        PANEL_TOP + Inches(0.10),
        LEFT_W - Inches(0.20),
        Inches(0.22),
        size=10,
        bold=True,
        color=CGI_RED,
    )

    tree = [
        ("output_30day/", False, 0),
        ("  state.db  (31 tables, WAL mode)", True, 1),
        ("  run_log.txt  (30 advance runs)", True, 1),
        ("  deltas/", False, 0),
        ("    2026-04-01/  \u2190 sample day", False, 1),
        ("      inserts/", False, 1),
        ("        TRANSACTION.csv       574  + .json", True, 2),
        ("        AUTHORIZATION.csv     634  + .json", True, 2),
        ("        STATEMENT.csv          21  + .json", True, 2),
        ("        COLLECTION_CASE.csv    27  + .json", True, 2),
        ("        FRAUD_ALERT.csv         2  + .json", True, 2),
        ("        SCORE_RECORD.csv      494  + .json", True, 2),
        ("      updates/", False, 1),
        ("        ACCOUNT.csv           494  + .json", True, 2),
        ("        CARD.csv              465  + .json", True, 2),
        ("        FRAUD_ALERT.csv         2  + .json", True, 2),
        ("    2026-04-02  \u2192  2026-04-30  \u2026", False, 1),
    ]

    ty = PANEL_TOP + Inches(0.40)
    for text, is_file, _ind in tree:
        color = CGI_BODY if is_file else CGI_BLACK
        bold = not is_file
        txt(
            slide,
            text,
            MARGIN_L + Inches(0.14),
            ty,
            LEFT_W - Inches(0.22),
            Inches(0.24),
            size=10,
            bold=bold,
            color=color,
        )
        ty += Inches(0.257)

    # ── Right: ACCOUNT + CUSTOMER mini-tables ───────────────────────────────
    acc_cols = ["acct_id", "account_number", "credit_limit", "balance", "status", "days_del"]
    acc_w0 = [Inches(0.56), Inches(1.90), Inches(1.08)]
    _fill = RIGHT_W - sum(acc_w0) - Inches(1.08) - Inches(1.48) - Inches(0.60)
    acc_widths = acc_w0 + [Inches(1.08), Inches(1.48), Inches(0.60)]

    acc_rows = [
        [1, "...8375", "$11,241", "$5,724", "ACTIVE", 0],
        [2, "...9029", "$14,489", "$7,422", "ACTIVE", 0],
        [3, "...2117", "$39,303", "$2,473", "DELINQUENT", 45],
        [4, "...6393", "$60,226", "$31,406", "DELINQUENT", 1],
        [5, "...6865", "$54,105", "$14,878", "ACTIVE", 0],
    ]

    bottom = _v4_mini_table(
        slide,
        "ACCOUNT  (500 rows  |  updated every advance day)",
        acc_cols,
        acc_rows,
        RIGHT_X,
        PANEL_TOP,
        acc_widths,
        CGI_RED,
    )

    cust_cols = ["cust_id", "first_name", "last_name", "city", "prov", "lang"]
    cust_widths = [
        Inches(0.56),
        Inches(1.28),
        Inches(1.28),
        Inches(1.68),
        Inches(0.56),
        Inches(0.56),
    ]

    cust_rows = [
        [1, "Benjamin", "Harris", "East Jenny", "PE", "FR"],
        [2, "Brian", "Simmons", "Port Leahmouth", "NS", "EN"],
        [3, "Mark", "Davis", "North Jeffberg", "PE", "FR"],
        [4, "Dawn", "Perry", "Jacobsshire", "BC", "EN"],
        [5, "William", "York", "Kimberlybury", "ON", "EN"],
    ]

    _v4_mini_table(
        slide,
        "CUSTOMER  (500 rows  \u2014  SIN synthetic, no real PII)",
        cust_cols,
        cust_rows,
        RIGHT_X,
        bottom + Inches(0.14),
        cust_widths,
        CGI_BLUE,
    )

    add_footer(slide, 6)


# ---------------------------------------------------------------------------
# V4 SLIDE 7  -  Sample Data: Transaction & Authorization  (30-day actuals)
# ---------------------------------------------------------------------------


def slide_v4_07_txn_auth(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Sample Data  \u2014  Transaction & Authorization",
        "30-day run  |  ~581 transactions/day  |  ~633 authorizations/day  |  approved + declined",
    )

    # TRANSACTION
    tx_cols = ["txn_id", "acct_id", "date", "amount", "type", "description", "status"]
    tx_rows = [
        [1, 2, "2026-04-01 13:52", "$8.02", "INTEREST", "COSTCO WHOLESALE #1409", "posted"],
        [4, 6, "2026-04-01 04:48", "$449.76", "PURCHASE", "APPLE.COM/BILL", "posted"],
        [5, 7, "2026-04-01 15:32", "$208.38", "PAYMENT", "SHELL OIL 23434", "posted"],
        [8, 7, "2026-04-01 19:57", "$263.04", "PURCHASE", "TIM HORTONS #1488", "posted"],
        [112, 93, "2026-04-01 17:55", "$7.92", "FEE", "METRO INC #8177", "posted"],
    ]
    tx_widths = [
        Inches(0.52),
        Inches(0.56),
        Inches(1.52),
        Inches(0.78),
        Inches(1.24),
        Inches(3.16),
        Inches(0.62),
    ]
    _v4_mini_table(
        slide,
        "TRANSACTION  (~581 new rows/day  |  17,420 total after 30 days)",
        tx_cols,
        tx_rows,
        MARGIN_L,
        Inches(1.56),
        tx_widths,
        CGI_RED,
    )

    # AUTHORIZATION
    auth_cols = [
        "auth_id",
        "acct_id",
        "amount",
        "response",
        "approval_code",
        "network",
        "channel",
        "decline_reason",
    ]
    auth_rows = [
        [3, 2, "$8.02", "00 Approved", "177852", "Interac", "telephone", "\u2014"],
        [5, 5, "$336.19", "00 Approved", "959001", "BankNet", "atm", "\u2014"],
        [6, 5, "$406.25", "00 Approved", "471027", "AmexNet", "telephone", "\u2014"],
        [1, 1, "$681.24", "51 Declined", "\u2014", "VisaNet", "mobile", "do_not_honour"],
        [4, 4, "$526.49", "54 Declined", "\u2014", "BankNet", "in_store", "insufficient_funds"],
    ]
    auth_widths = [
        Inches(0.56),
        Inches(0.56),
        Inches(0.78),
        Inches(0.92),
        Inches(0.96),
        Inches(0.82),
        Inches(0.86),
        Inches(1.60),
    ]
    _v4_mini_table(
        slide,
        "AUTHORIZATION  (~633 new rows/day  |  19,003 total  |  includes declined auths)",
        auth_cols,
        auth_rows,
        MARGIN_L,
        Inches(3.90),
        auth_widths,
        RGBColor(0x1A, 0x7A, 0x3C),
    )

    add_footer(slide, 7)


# ---------------------------------------------------------------------------
# V4 SLIDE 8  -  Sample Data: Dispute, Fraud Alert & Collection (30-day actuals)
# ---------------------------------------------------------------------------


def slide_v4_08_lifecycle(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "Sample Data  \u2014  Dispute, Fraud Alert & Collection Case",
        "Lifecycle entities progress through states each day  |  30-day run shown",
    )

    # DISPUTE
    disp_cols = ["dispute_id", "txn_id", "type", "amount", "opened", "status", "resolution"]
    disp_rows = [
        [1, 1802, "QUALITY", "$201.56", "2026-04-04", "RESOLVED", "APPROVED"],
        [2, 1954, "FRAUD", "$60.90", "2026-04-04", "WITHDRAWN", "\u2014"],
        [3, 2163, "SUBSCRIPTION", "$340.94", "2026-04-04", "RESOLVED", "PARTIAL"],
        [5, 5018, "FRAUD", "$585.61", "2026-04-09", "RESOLVED", "DENIED"],
        [8, 6559, "WRONG_AMOUNT", "$224.86", "2026-04-12", "RESOLVED", "DENIED"],
    ]
    disp_widths = [
        Inches(0.80),
        Inches(0.70),
        Inches(1.18),
        Inches(0.80),
        Inches(0.96),
        Inches(1.00),
        Inches(0.96),
    ]
    _v4_mini_table(
        slide,
        "DISPUTE  (35 total  |  18 days with new disputes  |  multiple resolution outcomes)",
        disp_cols,
        disp_rows,
        MARGIN_L,
        Inches(1.56),
        disp_widths,
        CGI_RED,
    )

    # FRAUD_ALERT
    fa_cols = [
        "alert_id",
        "acct_id",
        "source",
        "type",
        "risk_score",
        "status",
        "action",
        "resolved",
    ]
    fa_rows = [
        [1, 28, "3DS", "LARGE_TXN", 701, "FALSE_POSITIVE", "none", "2026-04-02"],
        [4, 33, "rule_engine", "GEO_ANOMALY", 736, "CONFIRMED", "block_card", "2026-04-04"],
        [5, 121, "rule_engine", "LARGE_TXN", 759, "CONFIRMED", "block_card", "2026-04-05"],
        [6, 481, "3DS", "HIGH_RISK_MCC", 908, "FALSE_POSITIVE", "none", "2026-04-16"],
        [8, 208, "RTD", "LARGE_TXN", 826, "CONFIRMED", "block_card", "2026-04-22"],
    ]
    fa_widths = [
        Inches(0.62),
        Inches(0.62),
        Inches(0.92),
        Inches(1.12),
        Inches(0.82),
        Inches(1.14),
        Inches(0.88),
        Inches(0.92),
    ]
    _v4_mini_table(
        slide,
        "FRAUD_ALERT  (11 total  |  5 confirmed, 6 false-positive  |  CONFIRMED triggers card block)",
        fa_cols,
        fa_rows,
        MARGIN_L,
        Inches(3.86),
        fa_widths,
        RGBColor(0xCC, 0x44, 0x00),
    )

    # COLLECTION_CASE
    cc_cols = ["case_id", "acct_id", "opened", "bucket", "past_due", "total_owed", "status"]
    cc_rows = [
        [1, 3, "2026-04-01", "B2", "$49.46", "$2,473", "PROMISE_TO_PAY"],
        [2, 42, "2026-04-01", "B2", "$696.89", "$34,844", "ACTIVE"],
        [3, 81, "2026-04-01", "B3", "$417.80", "$20,890", "PROMISE_TO_PAY"],
        [27, 493, "2026-04-01", "B2", "$263.87", "$13,194", "PROMISE_TO_PAY"],
    ]
    cc_widths = [
        Inches(0.62),
        Inches(0.62),
        Inches(0.94),
        Inches(0.62),
        Inches(0.80),
        Inches(0.86),
        Inches(1.28),
    ]
    _v4_mini_table(
        slide,
        "COLLECTION_CASE  (27 total  |  auto-opened at days_delinquent \u2265 30  |  buckets advance daily)",
        cc_cols,
        cc_rows,
        MARGIN_L,
        Inches(5.42),
        cc_widths,
        CGI_BLUE,
    )

    add_footer(slide, 8)


# ---------------------------------------------------------------------------
# V4 SLIDE 9  -  30-Day Run Summary
# ---------------------------------------------------------------------------


def slide_v4_09_run_summary(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "30-Day Run Summary",
        "init 2026-03-31 (Day 0)  +  advance \u00d730  |  ~0.35s per day  |  500 accounts",
    )

    # Three-column run log: days 1-10 / 11-20 / 21-30 (+init in col1)
    COL_W = (CONTENT_W - 2 * Inches(0.14)) // 3
    COL_G = Inches(0.14)
    TOP = Inches(1.56)
    ROW_H = Inches(0.245)
    HDR_H = Inches(0.285)

    # Sub-column widths within each COL_W
    sub_w = [Inches(1.28), Inches(0.75), Inches(0.68), Inches(0.68), Inches(0.48)]
    # Adjust last to fill: COL_W - sum(first 4)
    sub_w[4] = COL_W - sum(sub_w[:4])

    sub_h = ["Date", "Inserts", "Txns", "Updates", "Dur"]

    # Split run log: init+days1-10 / days11-20 / days21-30
    col_data = [
        _V4_RUN_LOG[0:11],  # init + days 1-10
        _V4_RUN_LOG[11:21],  # days 11-20
        _V4_RUN_LOG[21:31],  # days 21-30
    ]

    for ci, rows in enumerate(col_data):
        cx = MARGIN_L + ci * (COL_W + COL_G)
        # header bar
        rect(slide, cx, TOP, COL_W, HDR_H, CGI_RED)
        hx = cx
        for cw, ch in zip(sub_w, sub_h):
            txt(
                slide,
                ch,
                hx + Inches(0.04),
                TOP + Inches(0.04),
                cw - Inches(0.06),
                HDR_H - Inches(0.06),
                size=9,
                bold=True,
                color=CGI_WHITE,
            )
            hx += cw

        ry = TOP + HDR_H
        for ri, (date, ins, txns, upd, dur) in enumerate(rows):
            bg = CGI_LIGHT_GREY if ri % 2 == 0 else CGI_WHITE
            rx = cx
            for cw, val in zip(
                sub_w,
                [
                    date,
                    f"{ins:,}",
                    f"{txns:,}" if txns else "\u2014",
                    f"{upd:,}" if upd else "\u2014",
                    dur,
                ],
            ):
                rect(slide, rx, ry, cw, ROW_H, bg)
                txt(
                    slide,
                    str(val),
                    rx + Inches(0.04),
                    ry + Inches(0.03),
                    cw - Inches(0.06),
                    ROW_H - Inches(0.04),
                    size=9,
                    color=CGI_BODY,
                )
                rx += cw
            ry += ROW_H

    # ── Cumulative Final State ───────────────────────────────────────────────
    STATE_TOP = TOP + HDR_H + 11 * ROW_H + Inches(0.16)
    STATE_ROW = Inches(0.255)
    STATE_HDR = Inches(0.28)

    txt(
        slide,
        "CUMULATIVE STATE  (Day 30  \u2014  2026-04-30)",
        MARGIN_L,
        STATE_TOP,
        CONTENT_W,
        Inches(0.22),
        size=10,
        bold=True,
        color=CGI_RED,
    )

    STATE_TOP += Inches(0.26)

    # Three columns of 5 entities each
    groups = [_V4_FINAL_STATE[0:5], _V4_FINAL_STATE[5:10], _V4_FINAL_STATE[10:15]]
    for ci, group in enumerate(groups):
        cx = MARGIN_L + ci * (COL_W + COL_G)
        nw = COL_W - Inches(0.78)
        rw = Inches(0.78)

        rect(slide, cx, STATE_TOP, COL_W, STATE_HDR, CGI_BLACK)
        txt(
            slide,
            "Table",
            cx + Inches(0.06),
            STATE_TOP + Inches(0.04),
            nw,
            STATE_HDR - Inches(0.06),
            size=9,
            bold=True,
            color=CGI_WHITE,
        )
        txt(
            slide,
            "Rows",
            cx + nw,
            STATE_TOP + Inches(0.04),
            rw - Inches(0.06),
            STATE_HDR - Inches(0.06),
            size=9,
            bold=True,
            color=CGI_WHITE,
            align=PP_ALIGN.RIGHT,
        )

        ry = STATE_TOP + STATE_HDR
        for ri, (tbl, cnt) in enumerate(group):
            bg = CGI_LIGHT_GREY if ri % 2 == 0 else CGI_WHITE
            rect(slide, cx, ry, COL_W, STATE_ROW, bg)
            txt(
                slide,
                tbl,
                cx + Inches(0.06),
                ry + Inches(0.03),
                nw - Inches(0.06),
                STATE_ROW - Inches(0.04),
                size=9,
                bold=False,
                color=CGI_BODY,
            )
            txt(
                slide,
                f"{cnt:,}",
                cx + nw,
                ry + Inches(0.03),
                rw - Inches(0.08),
                STATE_ROW - Inches(0.04),
                size=9,
                bold=False,
                color=CGI_BODY,
                align=PP_ALIGN.RIGHT,
            )
            ry += STATE_ROW

    # Summary line
    SUMMARY_Y = STATE_TOP + STATE_HDR + 5 * STATE_ROW + Inches(0.10)
    txt(
        slide,
        "31 runs (init + 30 advance)  |  39,221 total inserts  |  29,736 total updates  "
        "|  ~0.35s/day  |  output_30day/state.db",
        MARGIN_L,
        SUMMARY_Y,
        CONTENT_W,
        Inches(0.26),
        size=10,
        bold=False,
        color=CGI_GREY,
    )

    add_footer(slide, 9)


# ---------------------------------------------------------------------------
# V4 SLIDE 10  -  CLI Reference  (actual commands from 30-day run)
# ---------------------------------------------------------------------------


def slide_v4_10_cli(prs):
    slide = blank_slide(prs)
    slide_header(
        slide,
        "CLI Reference",
        "All modes run from the terminal \u2014 no GUI required  |  actual commands from the 30-day run",
    )

    cmds = [
        (
            "init",
            "python -m tdgen_temporal.cli init --db output_30day/state.db --date 2026-03-31",
            "Seed the Day 0 population. Run once. Creates output_30day/state.db with 500 accounts,\n"
            "500 customers, 500 cards, 200 merchants, and all reference tables.  (1,733 inserts, ~0.4s)",
        ),
        (
            "advance",
            "python -m tdgen_temporal.cli advance --db output_30day/state.db --days 30 --output output_30day",
            "Advance the simulation by N days from the current state. Writes delta files to\n"
            "output_30day/deltas/YYYY-MM-DD/inserts/ and updates/ each day.  (~1,240 inserts/day, ~0.35s)",
        ),
        (
            "backfill",
            "python -m tdgen_temporal.cli backfill --from 2026-01-01 --to 2026-03-30 --db output_30day/state.db",
            "Generate a full date range of history in sequence. Useful for creating a rich\n"
            "historical dataset before handing off to QA teams.",
        ),
        (
            "status",
            "python -m tdgen_temporal.cli status --db output_30day/state.db",
            "Show current simulation date, total runs, and row counts for all entity tables.\n"
            "No side effects \u2014 safe to run at any time.",
        ),
    ]

    TOP = Inches(1.56)
    BOX_H = Inches(1.24)
    GAP = Inches(0.10)
    CMD_Y_OFF = Inches(0.38)
    CMD_H = Inches(0.28)

    ty = TOP
    for verb, cmd, desc in cmds:
        rect(slide, MARGIN_L, ty, CONTENT_W, BOX_H, CGI_LIGHT_GREY)
        rect(slide, MARGIN_L, ty, Inches(0.06), BOX_H, CGI_RED)

        txt(
            slide,
            verb,
            MARGIN_L + Inches(0.14),
            ty + Inches(0.08),
            Inches(0.90),
            Inches(0.24),
            size=13,
            bold=True,
            color=CGI_RED,
        )

        # Command box
        rect(
            slide,
            MARGIN_L + Inches(0.14),
            ty + CMD_Y_OFF,
            CONTENT_W - Inches(0.24),
            CMD_H,
            RGBColor(0x1A, 0x1A, 0x1A),
        )
        txt(
            slide,
            cmd,
            MARGIN_L + Inches(0.22),
            ty + CMD_Y_OFF + Inches(0.03),
            CONTENT_W - Inches(0.40),
            CMD_H - Inches(0.04),
            size=10,
            bold=False,
            color=RGBColor(0xAA, 0xFF, 0xAA),
        )

        txt(
            slide,
            desc,
            MARGIN_L + Inches(1.14),
            ty + Inches(0.06),
            CONTENT_W - Inches(1.24),
            Inches(0.30),
            size=11,
            color=CGI_BODY,
        )

        ty += BOX_H + GAP

    add_footer(slide, 10)


# ---------------------------------------------------------------------------
# V4 SLIDE 11  -  Next Steps  (reused from v2 slide_12_nextsteps)
# ---------------------------------------------------------------------------
# slide_12_nextsteps(prs) is called directly with footer override via wrapper


def slide_v4_11_nextsteps(prs):
    """Next Steps — v3 content, renumbered footer."""
    slide = blank_slide(prs)
    slide_header(
        slide, "Next Steps", "TDGen Temporal  |  Phase 1 complete  |  TD Bank / CGI engagement"
    )

    LEFT_W = Inches(5.90)
    RIGHT_X = MARGIN_L + LEFT_W + Inches(0.22)
    RIGHT_W = CONTENT_W - LEFT_W - Inches(0.22)

    TOP = Inches(1.56)

    def phase_box(x, y, w, title, items, accent):
        rect(slide, x, y, w, Inches(0.38), accent)
        txt(
            slide,
            title,
            x + Inches(0.14),
            y + Inches(0.06),
            w - Inches(0.20),
            Inches(0.28),
            size=14,
            bold=True,
            color=CGI_WHITE,
        )
        by = y + Inches(0.44)
        for item in items:
            txt(
                slide,
                "\u25cf  " + item,
                x + Inches(0.14),
                by,
                w - Inches(0.22),
                Inches(0.26),
                size=11,
                color=CGI_BODY,
            )
            by += Inches(0.30)
        return by

    phase1 = [
        "Day-over-day temporal progression \u2014 accounts, balances, delinquency",
        "6 entity state machines \u2014 ACCOUNT, CARD, DISPUTE, FRAUD_ALERT, CHARGEBACK, COLLECTION_CASE",
        "Daily transaction + authorization generation (~581 txns/day per 500 accounts)",
        "Statement generation on account cycle_day  (~18/day)",
        "Score record monthly refresh with drift based on account health",
        "INSERT and UPDATE delta files per day  (CSV + JSON)",
        "SQLite state store with WAL mode \u2014 persistent, queryable, restartable",
        "CLI: init / advance / backfill / status",
        "config/scenario.yaml \u2014 all rates configurable without code changes",
    ]

    phase2 = [
        "Validate with TD QA team \u2014 confirm rates match expected fraud lifecycle volumes",
        "New account opens \u2014 daily probabilistic seeding of new ACCOUNT + CUSTOMER + CARD",
        "MCP server wrapper \u2014 expose init/advance/status as MCP tools for Claude Code",
        "GitHub Actions workflow \u2014 scheduled nightly advance + delta archival",
        "Scenario presets \u2014 high-fraud day, payment cycle stress, chargeoff wave",
        "PostgreSQL provisioner \u2014 SQLAlchemy bulk insert for staging environment",
    ]

    phase_box(
        MARGIN_L, TOP, LEFT_W, "Phase 1  \u2014  Complete", phase1, RGBColor(0x1A, 0x7A, 0x3C)
    )
    phase_box(RIGHT_X, TOP, RIGHT_W, "Phase 2  \u2014  Planned", phase2, CGI_BLUE)

    add_footer(slide, 11)


# ---------------------------------------------------------------------------
# V4 SLIDE 12  -  CGI Back Page
# ---------------------------------------------------------------------------


def slide_v4_12_cgi_back(prs):
    slide = blank_slide(prs)

    grad_rect(
        slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, [(0.0, "E4002B"), (1.0, "8C0015")], 135
    )

    txt(
        slide,
        "Insights you can act on",
        Inches(0.90),
        Inches(1.20),
        Inches(11.53),
        Inches(0.80),
        size=32,
        bold=True,
        color=CGI_WHITE,
    )

    body = (
        "Founded in 1976, CGI is among the largest IT and business consulting services firms in the world.\n\n"
        "We are insights-driven and outcomes-focused to help accelerate returns on your investments. "
        "Across hundreds of locations worldwide, we provide comprehensive, scalable and sustainable "
        "IT and business consulting services that are informed globally and delivered locally."
    )
    txt(
        slide,
        body,
        Inches(0.90),
        Inches(2.20),
        Inches(11.53),
        Inches(3.80),
        size=16,
        bold=False,
        color=CGI_WHITE,
    )

    add_footer(slide, 12)


# ---------------------------------------------------------------------------
# Build V4
# ---------------------------------------------------------------------------


def build_v4_presentation(output_path: str = "docs/TDGen_Temporal_v4.pptx"):
    prs = new_prs()
    slide_v4_01_cover(prs)  # 1  Cover
    slide_02_what(prs)  # 2  Overview  (reused)
    slide_v4_03_scope_text(prs)  # 3  Scope: TSYS TS2
    slide_v4_04_state_machines(prs)  # 4  State Machines
    slide_v4_05_pipeline(prs)  # 5  9-Step Pipeline
    slide_v4_06_output_sample(prs)  # 6  Output Model + Account/Customer (COMBINED)
    slide_v4_07_txn_auth(prs)  # 7  Transaction & Authorization
    slide_v4_08_lifecycle(prs)  # 8  Dispute / Fraud Alert / Collection
    slide_v4_09_run_summary(prs)  # 9  30-Day Run Summary
    slide_v4_10_cli(prs)  # 10 CLI Reference
    slide_v4_11_nextsteps(prs)  # 11 Next Steps
    slide_v4_12_cgi_back(prs)  # 12 CGI Back Page

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    import sys

    if "--v4" in sys.argv:
        build_v4_presentation()
    else:
        build_presentation()
